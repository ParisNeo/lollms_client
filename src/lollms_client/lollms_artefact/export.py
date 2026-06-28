# lollms_discussion/_mixin_export.py
from __future__ import annotations

import io
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, TYPE_CHECKING

from ascii_colors import ASCIIColors, trace_exception
from .lollms_artefact import ArtefactType

def _parse_html_slides(html_content: str) -> List[Dict[str, Any]]:
    import pipmaster as pm
    try:
        pm.ensure_packages("beautifulsoup4")
        from bs4 import BeautifulSoup
    except Exception as e:
        ASCIIColors.warning(f"Could not load beautifulsoup4: {e}. Fallback to regex-based slide parser.")
        slides = []
        sections = re.findall(r'<section\s+class="slide[^"]*"[^>]*>(.*?)</section>', html_content, re.DOTALL | re.IGNORECASE)
        for idx, sec in enumerate(sections):
            notes_match = re.search(r'data-notes="([^"]*)"', sec, re.IGNORECASE)
            notes = notes_match.group(1).strip() if notes_match else ""

            title_match = re.search(r'<h[12][^>]*>(.*?)</h[12]>', sec, re.DOTALL | re.IGNORECASE)
            title = re.sub(r'<[^>]+>', '', title_match.group(1)).strip() if title_match else ""

            subtitle_match = re.search(r'<h3[^>]*>(.*?)</h3>', sec, re.DOTALL | re.IGNORECASE)
            subtitle = re.sub(r'<[^>]+>', '', subtitle_match.group(1)).strip() if subtitle_match else ""

            bullets = []
            for li in re.findall(r'<li[^>]*>(.*?)</li>', sec, re.DOTALL | re.IGNORECASE):
                bullets.append(re.sub(r'<[^>]+>', '', li).strip())

            img_match = re.search(r'<artefact_image\s+id=["\']([^"\']+)["\']', sec, re.IGNORECASE)
            image_id = img_match.group(1) if img_match else None

            slides.append({
                "index": idx,
                "title": title,
                "subtitle": subtitle,
                "bullets": [b for b in bullets if b],
                "notes": notes,
                "image_id": image_id
            })
        return slides

    soup = BeautifulSoup(html_content, "html.parser")
    slides = []

    sections = soup.find_all("section", class_="slide")
    if not sections:
        sections = soup.find_all("section")

    for idx, sec in enumerate(sections):
        notes = sec.get("data-notes", "").strip()

        title_el = sec.find(class_="slide-title") or sec.find(["h1", "h2"])
        title = title_el.get_text().strip() if title_el else ""

        subtitle_el = sec.find(class_="slide-subtitle") or sec.find("h3")
        subtitle = subtitle_el.get_text().strip() if subtitle_el else ""

        bullets = []
        bullet_list = sec.find(class_="slide-bullets") or sec.find(["ul", "ol"])
        if bullet_list:
            bullets = [li.get_text().strip() for li in bullet_list.find_all("li") if li.get_text().strip()]
        else:
            body_el = sec.find(class_="slide-body")
            if body_el:
                bullets = [p.get_text().strip() for p in body_el.find_all(["p", "div"]) if p.get_text().strip()]
            else:
                bullets = [p.get_text().strip() for p in sec.find_all("p") if p.get_text().strip()]

        if not title and (sec.find("h1") or sec.find("h2")):
            h = sec.find(["h1", "h2"])
            title = h.get_text().strip()

        # Find image anchor inside slide section
        img_match = re.search(r'<artefact_image\s+id=["\']([^"\']+)["\']', str(sec), re.IGNORECASE)
        image_id = img_match.group(1) if img_match else None

        slides.append({
            "index": idx,
            "title": title,
            "subtitle": subtitle,
            "bullets": bullets,
            "notes": notes,
            "image_id": image_id
        })

    return slides

class ExportMixin:
    """Provides methods to export artifacts in various formats."""
    
    def export_artifact(
        self,
        title: str,
        export_format: str,
        version: Optional[int] = None
    ) -> Tuple[bytes, str]:
        """
        Exports an artifact of the given title in the requested format.
        
        Returns:
            Tuple[bytes, str]: (binary_content, mime_type)
        """
        art = self.artefacts.get(title, version)
        if not art:
            raise ValueError(f"Artifact '{title}' not found.")
            
        fmt = export_format.lower().strip()
        atype = art.get("type", "document")
        
        if atype == "data":
            # Data-based exports
            if fmt == "csv":
                return self._export_data_as_csv(art), "text/csv"
            elif fmt in ("excel", "xlsx"):
                return self._export_data_as_excel(art), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            else:
                raise ValueError(f"Data artifacts can only be exported as 'csv' or 'excel'.")
        else:
            # Text-based exports
            if fmt == "markdown":
                return art["content"].encode("utf-8"), "text/markdown"
            elif fmt == "html":
                if atype == "presentation":
                    return art["content"].encode("utf-8"), "text/html"
                return self._export_text_as_html(art), "text/html"
            elif fmt == "pdf":
                content_stripped = art.get("content", "").strip()
                is_latex = (
                    art.get("language") in ("latex", "tex") or
                    content_stripped.startswith("\\documentclass") or
                    content_stripped.startswith("```latex") or
                    content_stripped.startswith("```tex")
                )
                if is_latex:
                    return self._export_latex_as_pdf(art), "application/pdf"
                return self._export_text_as_pdf(art), "application/pdf"
            elif fmt == "docx":
                return self._export_text_as_docx(art), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            elif fmt == "pptx":
                return self._export_text_as_pptx(art), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            elif fmt == "zip":
                return self._export_as_zip(art), "application/zip"
            else:
                raise ValueError(f"Unsupported export format '{export_format}' for text artifacts.")

    def _export_text_as_html(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("markdown")
        import markdown
        
        content = art.get("content", "")
        body_html = markdown.markdown(content)
        
        # Add basic aesthetic styling to the HTML template
        html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{art['title']}</title>
<style>
    body {{
        font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
        line-height: 1.6;
        max-width: 800px;
        margin: 40px auto;
        padding: 0 20px;
        color: #333;
    }}
    h1, h2, h3 {{ color: #111; margin-top: 24px; }}
    pre {{ background: #f4f4f4; padding: 12px; border-radius: 6px; overflow-x: auto; }}
    code {{ font-family: monospace; background: #f4f4f4; padding: 2px 4px; border-radius: 3px; }}
    table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
    th, td {{ border: 1px solid #ddd; padding: 8px 12px; text-align: left; }}
    th {{ background: #f5f5f5; }}
</style>
</head>
<body>
{body_html}
</body>
</html>
"""
        return html.encode("utf-8")

    def _export_latex_as_pdf(self, art: Dict[str, Any]) -> bytes:
        import shutil
        import re
        if not shutil.which("pdflatex"):
            raise RuntimeError(
                "LaTeX compilation engine ('pdflatex') was not found on your system.\n"
                "Please install a LaTeX distribution:\n"
                "  • Windows: MiKTeX (https://miktex.org/)\n"
                "  • macOS: MacTeX (https://www.tug.org/mactex/)\n"
                "  • Linux: TeX Live (sudo apt install texlive-latex-extra latexmk)"
            )

        import pipmaster as pm
        pm.ensure_packages("pdflatex")
        from pdflatex import PDFLaTeX

        content = art.get("content", "")

        # Strip any leading/trailing markdown code fences that the LLM might have used
        content_clean = content.strip()
        content_clean = re.sub(r'^```latex\s*|\s*```$', '', content_clean, flags=re.IGNORECASE).strip()
        content_clean = re.sub(r'^```tex\s*|\s*```$', '', content_clean, flags=re.IGNORECASE).strip()

        try:
            pdfl = PDFLaTeX.from_binarystring(content_clean.encode("utf-8"), "document")
            pdf_bytes, log, status = pdfl.create_pdf()
            if not pdf_bytes:
                raise ValueError(f"pdflatex returned empty output. Log:\n{log}")
            return pdf_bytes
        except Exception as e:
            raise RuntimeError(f"LaTeX Compilation Failed: {e}")

    def _export_text_as_pdf(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("reportlab")
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors

        content = art.get("content", "")
        buffer = io.BytesIO()
        atype = art.get("type", "document")
        is_presentation = atype == "presentation" or "article class=\"presentation\"" in content

        # Define custom background drawer to match the dark aesthetic "vibes" (#0b0f19)
        def draw_dark_background(canvas, document):
            canvas.saveState()
            canvas.setFillColor(colors.HexColor('#0b0f19'))
            canvas.rect(0, 0, document.pagesize[0], document.pagesize[1], fill=1, stroke=0)
            canvas.restoreState()

        if is_presentation:
            doc = SimpleDocTemplate(buffer, pagesize=landscape(letter), rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
        else:
            doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)

        styles = getSampleStyleSheet()
        story = []

        if is_presentation:
            slides_data = _parse_html_slides(content)
            if slides_data:
                h1_style = ParagraphStyle('H1Style', parent=styles['Heading2'], fontSize=24, leading=28, spaceBefore=12, spaceAfter=8, textColor=colors.HexColor('#f8fafc'))
                h2_style = ParagraphStyle('H2Style', parent=styles['Heading3'], fontSize=16, leading=20, spaceBefore=10, spaceAfter=6, textColor=colors.HexColor('#94a3b8'))
                normal_style = ParagraphStyle('NormalStyle', parent=styles['Normal'], fontSize=12, leading=16, spaceAfter=10, textColor=colors.HexColor('#cbd5e1'))
                notes_style = ParagraphStyle('NotesStyle', parent=normal_style, fontName='Helvetica-Oblique', fontSize=10, leading=13, textColor=colors.HexColor('#a0a0c0'))

                for s_idx, s_data in enumerate(slides_data):
                    if s_idx > 0:
                        story.append(PageBreak())
                    story.append(Paragraph(f"Slide {s_idx + 1}: {s_data['title']}", h1_style))
                    if s_data["subtitle"]:
                        story.append(Paragraph(s_data["subtitle"], h2_style))
                    story.append(Spacer(1, 15))
                    for bullet in s_data["bullets"]:
                        story.append(Paragraph(f"• {bullet}", normal_style))

                    # Embed active image artifacts if generated in this slide
                    if s_data.get("image_id"):
                        img_bytes = self._get_image_bytes_by_id(s_data["image_id"])
                        if img_bytes:
                            from reportlab.platypus import Image as RLImage
                            try:
                                rl_img = RLImage(io.BytesIO(img_bytes), width=240, height=135)
                                story.append(Spacer(1, 10))
                                story.append(rl_img)
                            except Exception as img_err:
                                ASCIIColors.warning(f"Failed to embed image in PDF: {img_err}")

                    if s_data["notes"]:
                        story.append(Spacer(1, 15))
                        story.append(Paragraph(f"<i>Speaker Notes:</i> {s_data['notes']}", notes_style))

                doc.build(story, onFirstPage=draw_dark_background, onLaterPages=draw_dark_background)
                return buffer.getvalue()

        # Title (for standard documents only)
        title_style = ParagraphStyle(
            'TitleStyle',
            parent=styles['Heading1'],
            fontSize=24,
            leading=28,
            textColor=colors.HexColor('#1e293b'),
            spaceAfter=20
        )
        story.append(Paragraph(art["title"], title_style))

        # Parse simple markdown into Paragraphs
        normal_style = styles['Normal']
        normal_style.fontSize = 10
        normal_style.leading = 14
        normal_style.spaceAfter = 8

        h1_style = ParagraphStyle('H1Style', parent=styles['Heading2'], fontSize=16, leading=20, spaceBefore=12, spaceAfter=6, textColor=colors.HexColor('#0f172a'))
        h2_style = ParagraphStyle('H2Style', parent=styles['Heading3'], fontSize=13, leading=16, spaceBefore=10, spaceAfter=4, textColor=colors.HexColor('#334155'))

        code_style = ParagraphStyle('CodeStyle', parent=normal_style, fontName='Courier', fontSize=9, leading=11, spaceAfter=8, backColor=colors.HexColor('#f1f5f9'), borderPadding=6)

        lines = content.splitlines()
        in_code = False
        code_block = []

        for line in lines:
            line_stripped = line.strip()
            if line_stripped.startswith("```"):
                if in_code:
                    # End code block
                    code_text = "<br/>".join(code_block)
                    story.append(Paragraph(code_text, code_style))
                    code_block = []
                    in_code = False
                else:
                    in_code = True
                continue

            if in_code:
                # Escape HTML tags inside code
                escaped = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                code_block.append(escaped)
                continue

            if line_stripped.startswith("# "):
                story.append(Paragraph(line_stripped[2:], h1_style))
            elif line_stripped.startswith("## "):
                story.append(Paragraph(line_stripped[3:], h1_style))
            elif line_stripped.startswith("### "):
                story.append(Paragraph(line_stripped[4:], h2_style))
            elif line_stripped:
                # Basic inline bold/italic parser (simple)
                parsed_line = line_stripped
                parsed_line = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', parsed_line)
                parsed_line = re.sub(r'\*(.*?)\*', r'<i>\1</i>', parsed_line)
                story.append(Paragraph(parsed_line, normal_style))
            else:
                story.append(Spacer(1, 6))

        doc.build(story)
        return buffer.getvalue()

    def _export_text_as_docx(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("python-docx")
        import docx
        
        content = art.get("content", "")
        doc = docx.Document()
        doc.add_heading(art["title"], level=0)
        
        lines = content.splitlines()
        in_code = False
        code_p = None
        
        for line in lines:
            line_stripped = line.strip()
            if line_stripped.startswith("```"):
                if in_code:
                    in_code = False
                    code_p = None
                else:
                    in_code = True
                    code_p = doc.add_paragraph()
                    code_p.style = 'Normal'
                continue
                
            if in_code:
                if code_p:
                    run = code_p.add_run(line + "\n")
                    run.font.name = 'Courier New'
                continue
                
            if line_stripped.startswith("# "):
                doc.add_heading(line_stripped[2:], level=1)
            elif line_stripped.startswith("## "):
                doc.add_heading(line_stripped[3:], level=2)
            elif line_stripped.startswith("### "):
                doc.add_heading(line_stripped[4:], level=3)
            elif line_stripped:
                doc.add_paragraph(line_stripped)
                
        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    def _export_text_as_pptx(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("python-pptx")
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        content = art.get("content", "")
        prs = Presentation()

        atype = art.get("type", "document")
        is_dark = 'data-theme="dark"' in content.lower()

        if is_dark:
            title_color = RGBColor(248, 250, 252)  # #f8fafc
            text_color = RGBColor(203, 213, 225)   # #cbd5e1
        else:
            title_color = RGBColor(15, 23, 42)     # #0f172a
            text_color = RGBColor(51, 65, 85)      # #334155

        def _apply_text_color(shape, color):
            if not shape.has_text_frame:
                return
            for p in shape.text_frame.paragraphs:
                p.font.color.rgb = color
                for r in p.runs:
                    r.font.color.rgb = color

        if atype == "presentation" or "article class=\"presentation\"" in content:
            slides_data = _parse_html_slides(content)
            if slides_data:
                # First slide is title slide
                first = slides_data[0]
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                
                if is_dark:
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(9, 9, 11)  # #09090b
                
                slide.shapes.title.text = first["title"] or art["title"]
                _apply_text_color(slide.shapes.title, title_color)
                
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = first["subtitle"] or "Exported from LoLLMS"
                    _apply_text_color(slide.placeholders[1], text_color)
                if first["notes"] and slide.notes_slide:
                    slide.notes_slide.notes_text_frame.text = first["notes"]

                # Remaining slides
                slide_layout_content = prs.slide_layouts[1]
                for s_data in slides_data[1:]:
                    slide = prs.slides.add_slide(slide_layout_content)
                    
                    if is_dark:
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = RGBColor(9, 9, 11)  # #09090b
                    
                    slide.shapes.title.text = s_data["title"]
                    _apply_text_color(slide.shapes.title, title_color)
                    
                    tf = slide.placeholders[1].text_frame
                    
                    first_para = True
                    if s_data["subtitle"]:
                        tf.paragraphs[0].text = s_data["subtitle"]
                        _apply_text_color(slide.placeholders[1], text_color)
                        first_para = False
                    
                    for bullet in s_data["bullets"]:
                        # Clean bullet text: normalize spaces and strip newlines to prevent text wrapping on next line
                        bullet_clean = " ".join(bullet.split())
                        # Remove leading special bullet symbols to prevent duplicate bullets in PowerPoint
                        bullet_clean = re.sub(r'^[▶◆●■•\*\-\s\d\.\)]+', '', bullet_clean).strip()
                        if not bullet_clean:
                            continue
                        
                        if first_para:
                            p = tf.paragraphs[0]
                            p.text = bullet_clean
                            p.level = 0
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet_clean
                            p.level = 0
                        
                        # Apply custom text color inside the paragraph runs
                        p.font.color.rgb = text_color
                        for r in p.runs:
                            r.font.color.rgb = text_color
                            
                    # Add Image if found on the slide
                    if s_data.get("image_id"):
                        # Resize text placeholder to left half safely by preserving its inherited coordinates
                        tf_shape = slide.placeholders[1]
                        orig_left = tf_shape.left
                        orig_top = tf_shape.top
                        orig_height = tf_shape.height
                        
                        tf_shape.left = orig_left
                        tf_shape.top = orig_top
                        tf_shape.width = Inches(5.0)
                        tf_shape.height = orig_height
                        
                        img_bytes = self._get_image_bytes_by_id(s_data["image_id"])
                        if img_bytes:
                            slide.shapes.add_picture(
                                io.BytesIO(img_bytes),
                                left=Inches(6.2),
                                top=orig_top,
                                width=Inches(3.5)
                            )

                    if s_data["notes"] and slide.notes_slide:
                        slide.notes_slide.notes_text_frame.text = s_data["notes"]

                buffer = io.BytesIO()
                prs.save(buffer)
                return buffer.getvalue()

        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = art["title"]
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Exported from LoLLMS Multimodal Viewer"

        # Split content into slides on h1/h2 headings
        lines = content.splitlines()
        slide_layout_content = prs.slide_layouts[1]
        current_slide = None
        tf = None

        for line in lines:
            line_stripped = line.strip()
            if line_stripped.startswith("# ") or line_stripped.startswith("## "):
                title_text = line_stripped.split(" ", 1)[1]
                current_slide = prs.slides.add_slide(slide_layout_content)
                current_slide.shapes.title.text = title_text
                tf = current_slide.placeholders[1].text_frame
            elif line_stripped:
                if tf:
                    p = tf.add_paragraph()
                    p.text = line_stripped
                    p.level = 0

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _export_data_as_csv(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("pandas")
        import pandas as pd

        ext = art.get("file_ext", ".csv")
        title = art["title"]
        version = art.get("version", 1)
        workspace_dir = Path("./data_workspace")
        try:
            from lollms_client.app.server import APP_WORKSPACE_DIR
            if APP_WORKSPACE_DIR is not None:
                workspace_dir = APP_WORKSPACE_DIR
        except ImportError:
            pass
        file_path = workspace_dir / f"{title}_v{version}{ext}"

        if not file_path.exists():
            raise FileNotFoundError(f"Original dataset file not found at {file_path}")
            
        try:
            if ext in (".db", ".sqlite", ".sqlite3"):
                import sqlite3
                conn = sqlite3.connect(str(file_path))
                cursor = conn.cursor()
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = [row[0] for row in cursor.fetchall() if row[0] != "sqlite_sequence"]
                if not tables:
                    raise ValueError("No tables found in SQLite DB.")
                df = pd.read_sql_query(f"SELECT * FROM {tables[0]};", conn)
                conn.close()
            elif ext in (".xlsx", ".xls"):
                df = pd.read_excel(str(file_path))
            else:
                sep = ";" if ext == ".csv" and ";" in file_path.read_text(encoding="utf-8", errors="ignore").splitlines()[0] else ","
                df = pd.read_csv(str(file_path), sep=sep)
                
            out = io.StringIO()
            df.to_csv(out, index=False)
            return out.getvalue().encode("utf-8")
        except Exception as e:
            trace_exception(e)
            raise RuntimeError(f"CSV export failed: {e}")

    def _export_data_as_excel(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages(["pandas", "openpyxl"])
        import pandas as pd

        ext = art.get("file_ext", ".csv")
        title = art["title"]
        version = art.get("version", 1)
        workspace_dir = Path("./data_workspace")
        try:
            from lollms_client.app.server import APP_WORKSPACE_DIR
            if APP_WORKSPACE_DIR is not None:
                workspace_dir = APP_WORKSPACE_DIR
        except ImportError:
            pass
        file_path = workspace_dir / f"{title}_v{version}{ext}"

        if not file_path.exists():
            raise FileNotFoundError(f"Original dataset file not found at {file_path}")
            
        try:
            buffer = io.BytesIO()
            with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
                if ext in (".db", ".sqlite", ".sqlite3"):
                    import sqlite3
                    conn = sqlite3.connect(str(file_path))
                    cursor = conn.cursor()
                    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                    tables = [row[0] for row in cursor.fetchall()]
                    for table in tables:
                        df = pd.read_sql_query(f"SELECT * FROM {table};", conn)
                        df.to_excel(writer, sheet_name=table[:31], index=False)
                    conn.close()
                elif ext in (".xlsx", ".xls"):
                    xl = pd.ExcelFile(str(file_path))
                    for sheet in xl.sheet_names:
                        df = pd.read_excel(str(file_path), sheet_name=sheet)
                        df.to_excel(writer, sheet_name=sheet[:31], index=False)
                else:
                    sep = ";" if ext == ".csv" and ";" in file_path.read_text(encoding="utf-8", errors="ignore").splitlines()[0] else ","
                    df = pd.read_csv(str(file_path), sep=sep)
                    df.to_excel(writer, sheet_name="Data", index=False)
                    
            return buffer.getvalue()
        except Exception as e:
            trace_exception(e)
            raise RuntimeError(f"Excel export failed: {e}")

    def _get_image_bytes_by_id(self, image_id: str) -> Optional[bytes]:
        import base64
        parts = image_id.split("::")
        if len(parts) < 2:
            return None
        try:
            img_index = int(parts[-1])
        except ValueError:
            return None
        title = "::".join(parts[:-1])

        # Check both main and companion
        main_art = self.artefacts.get(title)
        comp_art = self.artefacts.get(f"{title}::images")

        target_art = None
        target_index = None

        if comp_art and img_index < len(comp_art.get("images", [])):
            target_art = comp_art
            target_index = img_index
        elif main_art and img_index < len(main_art.get("images", [])):
            target_art = main_art
            target_index = img_index

        if target_art:
            imgs = target_art.get("images", [])
            if 0 <= target_index < len(imgs):
                b64_str = imgs[target_index]
                if b64_str:
                    if ";base64," in b64_str:
                        b64_str = b64_str.split(";base64,")[1]
                    return base64.b64decode(b64_str)
        return None

    def _export_as_zip(self, art: Dict[str, Any]) -> bytes:
        import zipfile
        import io
        import base64

        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as z:
            content = art.get("content", "")
            title = art["title"]

            # 1. Parse and extract all embedded image anchors
            img_pattern = re.compile(r'<artefact_image\s+id=["\']([^"\']+)["\']', re.IGNORECASE)
            image_ids = img_pattern.findall(content)

            # Track replaced image paths to update HTML refs
            replacements = {}
            for img_id in image_ids:
                if "::" in img_id:
                    img_bytes = self._get_image_bytes_by_id(img_id)
                    if img_bytes:
                        parts = img_id.split("::")
                        img_index = int(parts[-1])
                        img_title = "::".join(parts[:-1])
                        # Save image to zip under images/
                        safe_filename = re.sub(r'[^a-zA-Z0-9_.-]', '_', f"{img_title}_{img_index}.png")
                        zip_image_path = f"images/{safe_filename}"
                        z.writestr(zip_image_path, img_bytes)
                        replacements[img_id] = zip_image_path

            # 2. Update all <artefact_image> XML tags with relative <img> links in the zip's HTML index
            updated_content = content
            for img_id, rel_path in replacements.items():
                tag_pattern = re.compile(rf'<artefact_image\s+id=["\']{re.escape(img_id)}["\']\s*(?:\/>|>)', re.IGNORECASE)
                updated_content = tag_pattern.sub(f'<img src="{rel_path}" style="width:100%; height:100%; object-fit:cover;" />', updated_content)

            # Write the main index file to the zip
            index_filename = self._get_filename_with_ext(title, art.get('type'), art.get('language'), art.get('file_ext'))
            z.writestr(index_filename, updated_content.encode("utf-8"))

            # 3. Export ALL active workspace artifacts (code, tools, notes, pages)
            # This ensures that complex multi-file applications (e.g. index.html + style.css + app.js + main.py)
            # are completely bundled together into the root of the ZIP archive.
            all_active = self.artefacts.list(active_only=True)
            for item in all_active:
                item_title = item["title"]
                if item_title != title and not item_title.endswith("::images"):
                    item_content = item.get("content", "")
                    if item_content:
                        # Re-write and resolve any image anchors inside sibling artifacts as well
                        sibling_images = img_pattern.findall(item_content)
                        sibling_replacements = {}
                        for s_img_id in sibling_images:
                            if "::" in s_img_id:
                                s_img_bytes = self._get_image_bytes_by_id(s_img_id)
                                if s_img_bytes:
                                    s_parts = s_img_id.split("::")
                                    s_img_index = int(s_parts[-1])
                                    s_img_title = "::".join(s_parts[:-1])
                                    s_safe_filename = re.sub(r'[^a-zA-Z0-9_.-]', '_', f"{s_img_title}_{s_img_index}.png")
                                    s_zip_image_path = f"images/{s_safe_filename}"
                                    try:
                                        # Only write the image if not already written to zip
                                        z.writestr(s_zip_image_path, s_img_bytes)
                                    except UserWarning:
                                        pass # Duplicate file inside zip is fine
                                    sibling_replacements[s_img_id] = s_zip_image_path

                        sibling_updated_content = item_content
                        for s_img_id, s_rel_path in sibling_replacements.items():
                            s_tag_pattern = re.compile(rf'<artefact_image\s+id=["\']{re.escape(s_img_id)}["\']\s*(?:\/>|>)', re.IGNORECASE)
                            sibling_updated_content = s_tag_pattern.sub(f'<img src="{s_rel_path}" style="width:100%; height:100%; object-fit:cover;" />', sibling_updated_content)

                        # Save clean file to the zip
                        sibling_filename = self._get_filename_with_ext(item_title, item.get('type'), item.get('language'), item.get('file_ext'))
                        z.writestr(sibling_filename, sibling_updated_content.encode("utf-8"))

            # 4. Find and bundle all active datasets (CSVs, Excel files) in the workspace
            try:
                from lollms_client.app.server import APP_WORKSPACE_DIR
                if APP_WORKSPACE_DIR and APP_WORKSPACE_DIR.exists():
                    for active_art in self.artefacts.list(active_only=True):
                        if active_art.get("type") == "data":
                            d_title = active_art["title"]
                            ext = active_art.get("file_ext", ".csv")
                            version = active_art.get("version", 1)
                            
                            # Trim extension suffix from the title to prevent double-extension bugs during export
                            base_title = d_title
                            if base_title.lower().endswith(ext.lower()):
                                base_title = base_title[:-len(ext)]
                            
                            # We save both the versioned and unversioned names inside the zip so the code works under any reference style!
                            for name in (f"{base_title}{ext}", f"{base_title}_v{version}{ext}"):
                                file_path = APP_WORKSPACE_DIR / name
                                if file_path.exists():
                                    try:
                                        z.write(str(file_path), name)
                                    except UserWarning:
                                        pass
            except Exception:
                pass

        return buffer.getvalue()
    def _get_image_bytes_by_id(self, image_id: str) -> Optional[bytes]:
        import base64
        parts = image_id.split("::")
        if len(parts) < 2:
            return None
        try:
            img_index = int(parts[-1])
        except ValueError:
            return None
        title = "::".join(parts[:-1])

        # Check both main and companion
        main_art = self.artefacts.get(title)
        comp_art = self.artefacts.get(f"{title}::images")

        target_art = None
        target_index = None

        if comp_art and img_index < len(comp_art.get("images", [])):
            target_art = comp_art
            target_index = img_index
        elif main_art and img_index < len(main_art.get("images", [])):
            target_art = main_art
            target_index = img_index

        if target_art:
            imgs = target_art.get("images", [])
            if 0 <= target_index < len(imgs):
                b64_str = imgs[target_index]
                if b64_str:
                    if ";base64," in b64_str:
                        b64_str = b64_str.split(";base64,")[1]
                    return base64.b64decode(b64_str)
        return None

    def _export_as_zip(self, art: Dict[str, Any]) -> bytes:
        import zipfile
        import io
        import base64

        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as z:
            content = art.get("content", "")
            title = art["title"]

            # 1. Parse and extract all embedded image anchors
            img_pattern = re.compile(r'<artefact_image\s+id=["\']([^"\']+)["\']', re.IGNORECASE)
            image_ids = img_pattern.findall(content)

            # Track replaced image paths to update HTML refs
            replacements = {}
            for img_id in image_ids:
                if "::" in img_id:
                    img_bytes = self._get_image_bytes_by_id(img_id)
                    if img_bytes:
                        parts = img_id.split("::")
                        img_index = int(parts[-1])
                        img_title = "::".join(parts[:-1])
                        # Save image to zip under images/
                        safe_filename = re.sub(r'[^a-zA-Z0-9_.-]', '_', f"{img_title}_{img_index}.png")
                        zip_image_path = f"images/{safe_filename}"
                        z.writestr(zip_image_path, img_bytes)
                        replacements[img_id] = zip_image_path

            # 2. Update all <artefact_image> XML tags with relative <img> links in the zip's HTML index
            updated_content = content
            for img_id, rel_path in replacements.items():
                tag_pattern = re.compile(rf'<artefact_image\s+id=["\']{re.escape(img_id)}["\']\s*(?:\/>|>)', re.IGNORECASE)
                updated_content = tag_pattern.sub(f'<img src="{rel_path}" style="width:100%; height:100%; object-fit:cover;" />', updated_content)

            # Write the main index file to the zip
            index_filename = f"{title}.html" if not title.endswith(".html") else title
            z.writestr(index_filename, updated_content.encode("utf-8"))

            # 3. Search and bundle any sister/companion files (e.g., CSS/JS artifacts)
            related = self.artefacts.list()
            for r in related:
                r_title = r["title"]
                if r_title != title and (r_title.startswith(title.split(".")[0]) or r_title == f"{title}::images"):
                    if r_title == f"{title}::images":
                        continue  # Images already processed individually
                    r_content = r.get("content", "")
                    if r_content:
                        z.writestr(r_title, r_content.encode("utf-8"))

        return buffer.getvalue()

    def _get_image_bytes_by_id(self, image_id: str) -> Optional[bytes]:
        import base64
        parts = image_id.split("::")
        if len(parts) < 2:
            return None
        try:
            img_index = int(parts[-1])
        except ValueError:
            return None
        title = "::".join(parts[:-1])

        # Check both main and companion
        main_art = self.artefacts.get(title)
        comp_art = self.artefacts.get(f"{title}::images")

        target_art = None
        target_index = None

        if comp_art and img_index < len(comp_art.get("images", [])):
            target_art = comp_art
            target_index = img_index
        elif main_art and img_index < len(main_art.get("images", [])):
            target_art = main_art
            target_index = img_index

        if target_art:
            imgs = target_art.get("images", [])
            if 0 <= target_index < len(imgs):
                b64_str = imgs[target_index]
                if b64_str:
                    if ";base64," in b64_str:
                        b64_str = b64_str.split(";base64,")[1]
                    return base64.b64decode(b64_str)
        return None

    def _export_as_zip(self, art: Dict[str, Any]) -> bytes:
        import zipfile
        import io
        import base64

        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as z:
            content = art.get("content", "")
            title = art["title"]

            # 1. Parse and extract all embedded image anchors
            img_pattern = re.compile(r'<artefact_image\s+id=["\']([^"\']+)["\']', re.IGNORECASE)
            image_ids = img_pattern.findall(content)

            # Track replaced image paths to update HTML refs
            replacements = {}
            for img_id in image_ids:
                if "::" in img_id:
                    img_bytes = self._get_image_bytes_by_id(img_id)
                    if img_bytes:
                        parts = img_id.split("::")
                        img_index = int(parts[-1])
                        img_title = "::".join(parts[:-1])
                        # Save image to zip under images/
                        safe_filename = re.sub(r'[^a-zA-Z0-9_.-]', '_', f"{img_title}_{img_index}.png")
                        zip_image_path = f"images/{safe_filename}"
                        z.writestr(zip_image_path, img_bytes)
                        replacements[img_id] = zip_image_path

            # 2. Update all <artefact_image> XML tags with relative <img> links in the zip's HTML index
            updated_content = content
            for img_id, rel_path in replacements.items():
                tag_pattern = re.compile(rf'<artefact_image\s+id=["\']{re.escape(img_id)}["\']\s*(?:\/>|>)', re.IGNORECASE)
                updated_content = tag_pattern.sub(f'<img src="{rel_path}" style="width:100%; height:100%; object-fit:cover;" />', updated_content)

            # Write the main index file to the zip
            index_filename = f"{title}.html" if not title.endswith(".html") else title
            z.writestr(index_filename, updated_content.encode("utf-8"))

            # 3. Search and bundle any sister/companion files (e.g., CSS/JS artifacts)
            related = self.artefacts.list()
            for r in related:
                r_title = r["title"]
                if r_title != title and (r_title.startswith(title.split(".")[0]) or r_title == f"{title}::images"):
                    if r_title == f"{title}::images":
                        continue  # Images already processed individually
                    r_content = r.get("content", "")
                    if r_content:
                        z.writestr(r_title, r_content.encode("utf-8"))

        return buffer.getvalue()
