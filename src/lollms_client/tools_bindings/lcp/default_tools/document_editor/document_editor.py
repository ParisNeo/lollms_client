import os
import re
import json
import getpass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    import fitz
except ImportError:
    fitz = None

TOOL_LIBRARY_NAME = "Document Editor"
TOOL_LIBRARY_DESC = "Surgically edits, annotates (batch capable), highlights, copies-pastes between documents, deletes pages, and replaces text globally in PDF, DOCX, and PPTX files."
TOOL_LIBRARY_ICON = "📝"

_DEFAULT_COMMENTER = getpass.getuser()

_LIGATURE_MAP = {
    '\ufb00': 'ff',
    '\ufb01': 'fi',
    '\ufb02': 'fl',
    '\ufb03': 'ffi',
    '\ufb04': 'ffl',
    '\ufb05': 'ft',
    '\ufb06': 'st',
}

_LIGATURE_TRANS = str.maketrans(_LIGATURE_MAP)


def init_tools_library(config: dict = None) -> None:
    global fitz
    try:
        import pipmaster as pm
        pm.ensure_packages("pymupdf")
        pm.ensure_packages("python-docx")
        pm.ensure_packages("python-pptx")
    except Exception as e:
        import ascii_colors
        ascii_colors.ASCIIColors.warning(f"[Document Editor] Failed to ensure dependencies: {e}")

    if fitz is None:
        try:
            import fitz as _fitz
            fitz = _fitz
        except ImportError:
            import ascii_colors
            ascii_colors.ASCIIColors.error("[Document Editor] PyMuPDF (fitz) could not be imported. PDF annotation and editing will fail.")

def _get_output_path(file_name: str, suffix: str = "_edited") -> str:
    p = Path(file_name)
    stem = p.stem
    for existing_suffix in ["_annotated", "_edited", "_inserted", "_updated", "_removed", "_pasted", "_global_replace", "_page_deleted"]:
        if stem.endswith(existing_suffix):
            stem = stem[:-len(existing_suffix)]
            break
    return str(p.with_name(f"{stem}{suffix}{p.suffix}"))

def _check_file_exists(file_name: str) -> Optional[Dict[str, Any]]:
    if not Path(file_name).is_file():
        return {
            "success": False,
            "error": f"File '{file_name}' not found."
        }
    return None

def _parse_page_numbers(pages_str: Optional[str], total_pages: int) -> List[int]:
    if not pages_str or not pages_str.strip():
        return list(range(total_pages))
    
    pages = set()
    for part in str(pages_str).split(','):
        part = part.strip()
        if '-' in part:
            try:
                start, end = map(int, part.split('-'))
                pages.update(range(max(1, start) - 1, min(total_pages, end)))
            except ValueError:
                pass
        elif part.isdigit():
            p = int(part) - 1
            if 0 <= p < total_pages:
                pages.add(p)
    return sorted(list(pages))


def _normalize_text_for_search(text: str) -> str:
    text = text.translate(_LIGATURE_TRANS)
    text = text.replace('\n', ' ').replace('\r', ' ')
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def _extract_keywords(search_text: str, min_len: int = 4) -> List[str]:
    stop_words = {
        "the", "a", "an", "is", "are", "was", "were", "be", "been",
        "have", "has", "had", "do", "does", "did", "will", "would",
        "could", "should", "may", "might", "must", "shall", "can",
        "to", "of", "in", "for", "on", "with", "at", "by", "from",
        "as", "and", "or", "but", "not", "no", "if", "then", "so",
        "i", "you", "he", "she", "it", "we", "they",
    }
    normalized = _normalize_text_for_search(search_text)
    words = re.findall(r'[A-Za-z0-9+\-]+', normalized)
    return [w for w in words if len(w) >= min_len and w.lower() not in stop_words]


def _fuzzy_search_pdf_page(page, search_text: str, flags: int = 0) -> List:
    """
    Three-tier search for a text fragment on a PDF page using PyMuPDF.

    Returns a list of fitz.Rect objects (possibly empty).
    """
    if fitz is None:
        return []
    exact_hits = page.search_for(search_text, flags=flags)
    if exact_hits:
        return exact_hits

    norm_search = _normalize_text_for_search(search_text)
    if not norm_search:
        return []

    norm_hits = page.search_for(norm_search, flags=flags)
    if norm_hits:
        return norm_hits

    partial_fragments = page.search_for(norm_search[:80], flags=flags)
    if partial_fragments:
        return partial_fragments

    keywords = _extract_keywords(search_text, min_len=4)
    if not keywords:
        return []

    page_rect = page.rect
    keyword_rects: List = []
    for kw in keywords:
        kw_hits = page.search_for(kw, flags=flags)
        if kw_hits:
            keyword_rects.extend(kw_hits)

    if not keyword_rects:
        return []

    best_center: Optional[Tuple[float, float]] = None
    best_score = 0

    for anchor_rect in keyword_rects:
        center = (anchor_rect.x0 + anchor_rect.x1) / 2, (anchor_rect.y0 + anchor_rect.y1) / 2
        score = 0
        for other_rect in keyword_rects:
            oc = (other_rect.x0 + other_rect.x1) / 2, (other_rect.y0 + other_rect.y1) / 2
            dist = ((center[0] - oc[0]) ** 2 + (center[1] - oc[1]) ** 2) ** 0.5
            if dist < 200:
                score += 1
        if score > best_score:
            best_score = score
            best_center = center

    if best_center is None:
        return []

    cx, cy = best_center
    half_w = min(150, page_rect.width / 3)
    half_h = min(40, page_rect.height / 4)
    fuzzy_rect = page_rect & fitz.Rect(
        cx - half_w, cy - half_h, cx + half_w, cy + half_h
    )
    return [fuzzy_rect] if fuzzy_rect else []


def _fuzzy_find_in_paragraphs(paragraphs_text: List[str], search_text: str) -> List[int]:
    """
    Fuzzy search across a list of paragraph text strings.
    Returns indices of paragraphs that match.
    
    Tier 1: Exact substring match.
    Tier 2: Normalized match (whitespace/ligature insensitive).
    Tier 3: Keyword coverage match (>=60% of significant keywords present).
    """
    norm_search = _normalize_text_for_search(search_text)
    keywords = _extract_keywords(search_text, min_len=4)
    
    matched_indices: List[int] = []
    
    for idx, ptext in enumerate(paragraphs_text):
        if search_text in ptext:
            matched_indices.append(idx)
            continue
        
        norm_ptext = _normalize_text_for_search(ptext)
        if norm_search and norm_search in norm_ptext:
            matched_indices.append(idx)
            continue
        
        if keywords:
            present = sum(1 for kw in keywords if kw.lower() in norm_ptext.lower())
            if present / len(keywords) >= 0.6:
                matched_indices.append(idx)
                continue
    
    return matched_indices


def _apply_single_pdf_edit(page, operation: str, search_text: str, replacement_text: str, flags: int) -> int:
    """Applies a single edit operation to a PDF page. Returns number of matches."""
    matches = 0
    if operation == "insert":
        text_instances = _fuzzy_search_pdf_page(page, search_text, flags=flags)
        if text_instances:
            rect = text_instances[0]
            point = fitz.Point(rect.x1, rect.y0)
            page.insert_text(point, " " + replacement_text, fontsize=11, fontname="helv")
            matches += 1
    else:
        text_instances = _fuzzy_search_pdf_page(page, search_text, flags=flags)
        for inst in text_instances:
            if operation == "remove":
                page.add_redact_annot(inst, fill=(1, 1, 1))
            else:  # update
                page.add_redact_annot(inst, text=replacement_text, fill=(1, 1, 1), fontsize=11)
            matches += 1
        if matches > 0:
            page.apply_redactions()
    return matches


def _apply_single_pdf_annotation(page, annotation_type: str, search_text: str, comment: str, color: Tuple, commenter: str, flags: int) -> int:
    """Applies a single annotation to a PDF page. Returns number of annotations."""
    if fitz is None:
        return 0
    annot_count = 0
    text_instances = _fuzzy_search_pdf_page(page, search_text, flags=flags)

    for inst in text_instances:
        if annotation_type == "highlight":
            annot = page.add_highlight_annot(inst)
            annot.set_colors(stroke=color)
            annot.update()
        else:  # comment
            annot = page.add_text_annot(
                fitz.Point(inst.x0, inst.y0),
                comment,
                icon="Comment"
            )
            annot.set_info(title=commenter, content=comment)
            annot.update()
        annot_count += 1
    return annot_count


def tool_edit_document_text(
    file_name: str,
    operation: str,
    search_text: str,
    replacement_text: str = "",
    pages: str = "",
    match_case: bool = False,
    whole_word: bool = False
) -> Dict[str, Any]:
    """
    Surgically edits text content in a PDF, DOCX, or PPTX document.
    Uses fuzzy matching to locate the text and applies the specified operation.

    Args:
        file_name (str): The path to the document file (PDF, DOCX, PPTX).
        operation (str): The edit operation to perform. Options: "insert", "update", "remove".
        search_text (str): The exact text to search for. For "insert", this is the anchor text after which to insert.
        replacement_text (str, optional): The new text. Required for "update" and "insert". Ignored for "remove".
        pages (str, optional): Pages to apply the edit (PDF only). E.g., "1-3, 5". Empty means all pages.
        match_case (bool, optional): Whether to match the exact case. Defaults to False.
        whole_word (bool, optional): Whether to match whole words only. Defaults to False.
    """
    err = _check_file_exists(file_name)
    if err:
        return err

    if operation not in ("insert", "update", "remove"):
        return {"success": False, "error": "Invalid operation. Must be 'insert', 'update', or 'remove'."}
    
    if operation in ("update", "insert") and not replacement_text and operation != "remove":
        return {"success": False, "error": "replacement_text is required for 'update' and 'insert' operations."}

    output_path = _get_output_path(file_name, f"_{operation}ed")
    file_ext = Path(file_name).suffix.lower()

    try:
        if file_ext == ".pdf":
            if fitz is None:
                return {"success": False, "error": "PyMuPDF (fitz) is not installed. Cannot edit PDF files. Install it via 'pip install pymupdf'."}
            doc = fitz.open(file_name)
            target_pages = _parse_page_numbers(pages, len(doc))
            total_matches = 0
            
            flags = 0
            if not match_case:
                flags |= fitz.TEXT_DEHYPHENATE
            if whole_word:
                flags |= fitz.TEXT_FIND_WHOLEWORDS

            for page_num in target_pages:
                page = doc[page_num]
                total_matches += _apply_single_pdf_edit(page, operation, search_text, replacement_text, flags)

            if total_matches == 0:
                doc.close()
                return {"success": False, "error": f"Search text not found in the specified pages."}

            doc.save(output_path)
            doc.close()
            
        elif file_ext == ".docx":
            import docx
            doc = docx.Document(file_name)
            total_matches = 0
            
            paragraph_texts = [p.text for p in doc.paragraphs]
            matched_indices = _fuzzy_find_in_paragraphs(paragraph_texts, search_text)
            
            for idx in matched_indices:
                paragraph = doc.paragraphs[idx]
                if operation == "remove":
                    paragraph.clear()
                    total_matches += 1
                elif operation == "update":
                    paragraph.text = paragraph.text.replace(search_text, replacement_text)
                    if search_text not in paragraph.text:
                        norm_search = _normalize_text_for_search(search_text)
                        norm_ptext = _normalize_text_for_search(paragraph.text)
                        if norm_search in norm_ptext:
                            paragraph.text = norm_ptext.replace(norm_search, replacement_text)
                    total_matches += 1
                elif operation == "insert":
                    paragraph.add_run(" " + replacement_text)
                    total_matches += 1

            if total_matches == 0:
                return {"success": False, "error": "Search text not found in document."}

            doc.save(output_path)

        elif file_ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_name)
            total_matches = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        full_text = "".join(run.text for run in paragraph.runs)
                        matched = False
                        
                        if search_text in full_text:
                            matched = True
                        else:
                            norm_search = _normalize_text_for_search(search_text)
                            norm_full = _normalize_text_for_search(full_text)
                            if norm_search and norm_search in norm_full:
                                matched = True
                            else:
                                keywords = _extract_keywords(search_text, min_len=4)
                                if keywords:
                                    present = sum(1 for kw in keywords if kw.lower() in full_text.lower())
                                    if present / len(keywords) >= 0.6:
                                        matched = True
                        
                        if matched:
                            if operation == "remove":
                                for run in paragraph.runs:
                                    run.text = ""
                                total_matches += 1
                            elif operation == "update":
                                for run in paragraph.runs:
                                    if search_text in run.text:
                                        run.text = run.text.replace(search_text, replacement_text)
                                        total_matches += 1
                                    else:
                                        run.text = _normalize_text_for_search(run.text).replace(
                                            _normalize_text_for_search(search_text),
                                            replacement_text
                                        )
                                        if replacement_text in run.text:
                                            total_matches += 1
                            elif operation == "insert":
                                if paragraph.runs:
                                    paragraph.runs[-1].text += " " + replacement_text
                                    total_matches += 1

            if total_matches == 0:
                return {"success": False, "error": "Search text not found in presentation."}

            prs.save(output_path)
        else:
            return {"success": False, "error": f"Unsupported file extension: {file_ext}"}

        return {
            "success": True,
            "output": f"Successfully applied '{operation}' to {total_matches} instance(s). Saved to {output_path}",
            "file_name": output_path
        }

    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": f"Document editing failed: {str(e)}",
            "traceback": traceback.format_exc()
        }

def tool_annotate_document(
    file_name: str,
    annotation_type: str,
    search_text: str,
    comment: str = "",
    pages: str = "",
    highlight_color: str = "yellow",
    commenter_name: str = _DEFAULT_COMMENTER
) -> Dict[str, Any]:
    """
    Annotates a PDF or DOCX document by highlighting text or adding comments.

    Args:
        file_name (str): The path to the document file (PDF or DOCX).
        annotation_type (str): The type of annotation. Options: "comment", "highlight".
        search_text (str): The text to locate for annotation.
        comment (str, optional): The comment text. Required for "comment" type.
        pages (str, optional): Pages to apply the annotation (PDF only). E.g., "1-3, 5".
        highlight_color (str, optional): Color for highlighting. Options: "yellow", "red", "green", "blue".
        commenter_name (str, optional): The name of the user adding the comment. Defaults to current OS user.
    """
    err = _check_file_exists(file_name)
    if err:
        return err

    if annotation_type not in ("comment", "highlight"):
        return {"success": False, "error": "Invalid annotation_type. Must be 'comment' or 'highlight'."}

    if annotation_type == "comment" and not comment:
        return {"success": False, "error": "comment is required when annotation_type is 'comment'."}

    if not commenter_name:
        commenter_name = _DEFAULT_COMMENTER

    output_path = _get_output_path(file_name, "_annotated")
    file_ext = Path(file_name).suffix.lower()

    color_map = {
        "yellow": (1, 1, 0),
        "red": (1, 0, 0),
        "green": (0, 1, 0),
        "blue": (0, 0, 1)
    }
    rgb_color = color_map.get(highlight_color.lower(), (1, 1, 0))

    try:
        if file_ext == ".pdf":
            if fitz is None:
                return {"success": False, "error": "PyMuPDF (fitz) is not installed. Cannot annotate PDF files. Install it via 'pip install pymupdf'."}
            doc = fitz.open(file_name)
            target_pages = _parse_page_numbers(pages, len(doc))
            if not target_pages:
                doc.close()
                return {"success": False, "error": f"No valid pages parsed from '{pages}'. Document has {len(doc)} pages."}
            total_annot = 0

            for page_num in target_pages:
                page = doc[page_num]
                total_annot += _apply_single_pdf_annotation(page, annotation_type, search_text, comment, rgb_color, commenter_name, 0)

            if total_annot == 0:
                doc.close()
                return {"success": False, "error": f"Search text not found for annotation on page(s) {pages or 'all'}. The text may not exist or may be image-based (scanned PDF). Try a shorter, more specific search_text fragment."}

            doc.save(output_path)
            doc.close()

        elif file_ext == ".docx":
            import docx
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement

            doc = docx.Document(file_name)
            total_annot = 0

            paragraph_texts = [p.text for p in doc.paragraphs]
            matched_indices = _fuzzy_find_in_paragraphs(paragraph_texts, search_text)

            if not matched_indices:
                return {"success": False, "error": f"Search text not found in DOCX paragraphs. Tried exact, normalized, and keyword-coverage matching."}

            for idx in matched_indices:
                paragraph = doc.paragraphs[idx]
                if annotation_type == "highlight":
                    for run in paragraph.runs:
                        if search_text in run.text or _normalize_text_for_search(search_text) in _normalize_text_for_search(run.text):
                            rPr = run._element.get_or_create_rPr()
                            highlight = OxmlElement('w:highlight')
                            highlight.set(qn('w:val'), highlight_color)
                            rPr.append(highlight)
                            total_annot += 1
                    if total_annot == 0 and matched_indices:
                        for run in paragraph.runs:
                            rPr = run._element.get_or_create_rPr()
                            highlight = OxmlElement('w:highlight')
                            highlight.set(qn('w:val'), highlight_color)
                            rPr.append(highlight)
                            total_annot += 1
                else:
                    paragraph.add_run(f" [COMMENT by {commenter_name}: {comment}]")
                    total_annot += 1

            if total_annot == 0:
                return {"success": False, "error": f"Search text was found in paragraphs but no annotations could be applied. This may be due to empty runs or formatting constraints."}

            doc.save(output_path)
        else:
            return {"success": False, "error": f"Unsupported file extension: '{file_ext}'. Supported: .pdf, .docx"}

        return {
            "success": True,
            "output": f"Successfully added {total_annot} {annotation_type}(s). Saved to {output_path}",
            "file_name": output_path
        }

    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": f"Document annotation failed: {str(e)}",
            "traceback": traceback.format_exc()
        }

def tool_batch_annotate_document(
    file_name: str,
    batch_operations: str,
    pages: str = "",
    default_commenter_name: str = _DEFAULT_COMMENTER
) -> Dict[str, Any]:
    """
    Applies multiple annotations (highlights or comments) to a document in a single call.
    Useful for proofreading or reviewing a document with many distinct points.

    Args:
        file_name (str): The path to the document file (PDF or DOCX).
        batch_operations (str): A JSON string containing a list of operations.
            Each operation must be a dict with keys:
            - "annotation_type": "comment" or "highlight"
            - "search_text": The text to locate.
            - "comment": (Required for "comment" type) The comment text.
            - "highlight_color": (Optional for "highlight" type) "yellow", "red", "green", or "blue". Defaults to "yellow".
            - "commenter_name": (Optional) Overrides the default commenter name.
            Example: '[{"annotation_type": "highlight", "search_text": "error", "highlight_color": "red"}, {"annotation_type": "comment", "search_text": "syntax", "comment": "Fix syntax"}]'
        pages (str, optional): Pages to apply the annotations (PDF only). E.g., "1-3, 5". Empty means all pages.
        default_commenter_name (str, optional): The default name of the user adding comments. Defaults to current OS user.
    """
    err = _check_file_exists(file_name)
    if err:
        return err

    try:
        operations = json.loads(batch_operations)
        if not isinstance(operations, list) or not operations:
            return {"success": False, "error": "batch_operations must be a non-empty JSON array."}
    except json.JSONDecodeError as e:
        return {"success": False, "error": f"Invalid JSON in batch_operations: {e}"}

    output_path = _get_output_path(file_name, "_batch_annotated")
    file_ext = Path(file_name).suffix.lower()
    color_map = {
        "yellow": (1, 1, 0),
        "red": (1, 0, 0),
        "green": (0, 1, 0),
        "blue": (0, 0, 1)
    }

    total_applied = 0
    errors = []

    try:
        if file_ext == ".pdf":
            if fitz is None:
                return {"success": False, "error": "PyMuPDF (fitz) is not installed. Cannot batch annotate PDF files. Install it via 'pip install pymupdf'."}
            doc = fitz.open(file_name)
            target_pages = _parse_page_numbers(pages, len(doc))

            for op in operations:
                op_type = op.get("annotation_type")
                search_text = op.get("search_text", "")
                if not search_text:
                    errors.append("Missing search_text in operation.")
                    continue
                    
                color_name = op.get("highlight_color", "yellow").lower()
                rgb_color = color_map.get(color_name, (1, 1, 0))
                commenter = op.get("commenter_name", default_commenter_name)
                comment = op.get("comment", "")
                
                if op_type == "comment" and not comment:
                    errors.append(f"Missing comment for search_text: '{search_text}'")
                    continue
                
                op_count = 0
                for page_num in target_pages:
                    page = doc[page_num]
                    op_count += _apply_single_pdf_annotation(page, op_type, search_text, comment, rgb_color, commenter, 0)
                
                if op_count == 0:
                    errors.append(f"Text not found: '{search_text[:50]}...'")
                else:
                    total_applied += op_count
            
            if total_applied == 0:
                doc.close()
                return {"success": False, "error": "No annotations applied. " + " | ".join(errors)}
            
            doc.save(output_path)
            doc.close()

        elif file_ext == ".docx":
            import docx
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement

            doc = docx.Document(file_name)
            paragraph_texts = [p.text for p in doc.paragraphs]
            
            for op in operations:
                op_type = op.get("annotation_type")
                search_text = op.get("search_text", "")
                if not search_text:
                    errors.append("Missing search_text in operation.")
                    continue
                    
                color_name = op.get("highlight_color", "yellow").lower()
                commenter = op.get("commenter_name", default_commenter_name)
                comment = op.get("comment", "")
                
                if op_type == "comment" and not comment:
                    errors.append(f"Missing comment for search_text: '{search_text}'")
                    continue
                
                matched_indices = _fuzzy_find_in_paragraphs(paragraph_texts, search_text)
                if not matched_indices:
                    errors.append(f"Text not found: '{search_text[:50]}...'")
                    continue
                
                for idx in matched_indices:
                    paragraph = doc.paragraphs[idx]
                    if op_type == "highlight":
                        for run in paragraph.runs:
                            rPr = run._element.get_or_create_rPr()
                            highlight = OxmlElement('w:highlight')
                            highlight.set(qn('w:val'), color_name)
                            rPr.append(highlight)
                            total_applied += 1
                    else:
                        paragraph.add_run(f" [COMMENT by {commenter}: {comment}]")
                        total_applied += 1
            
            if total_applied == 0:
                return {"success": False, "error": "No annotations applied. " + " | ".join(errors)}
            
            doc.save(output_path)
        else:
            return {"success": False, "error": f"Unsupported file extension: {file_ext}"}

        result_msg = f"Successfully applied {total_applied} annotation(s). Saved to {output_path}"
        if errors:
            result_msg += f". Warnings: {' | '.join(errors)}"
            
        return {
            "success": True,
            "output": result_msg,
            "file_name": output_path
        }

    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": f"Batch annotation failed: {str(e)}",
            "traceback": traceback.format_exc()
        }

def tool_copy_paste_between_documents(
    source_file: str,
    target_file: str,
    source_search_text: str,
    target_anchor_text: str,
    source_pages: str = "",
    target_pages: str = ""
) -> Dict[str, Any]:
    """
    Copies a block of text from a source document and pastes it into a target document
    at a specified anchor location. Supports PDF documents.

    Args:
        source_file (str): The path to the source PDF document.
        target_file (str): The path to the target PDF document.
        source_search_text (str): The text in the source document to extract.
        target_anchor_text (str): The text in the target document after which to insert the copied text.
        source_pages (str, optional): Pages to search in the source PDF. E.g., "1-3, 5". Empty means all pages.
        target_pages (str, optional): Pages to search for the anchor in the target PDF. E.g., "1-3, 5". Empty means all pages.
    """
    err_src = _check_file_exists(source_file)
    if err_src:
        return err_src
    err_tgt = _check_file_exists(target_file)
    if err_tgt:
        return err_tgt

    src_ext = Path(source_file).suffix.lower()
    tgt_ext = Path(target_file).suffix.lower()

    if src_ext != ".pdf" or tgt_ext != ".pdf":
        return {"success": False, "error": "Copy-paste is currently only supported for PDF to PDF."}

    output_path = _get_output_path(target_file, "_pasted")

    try:
        if fitz is None:
            return {"success": False, "error": "PyMuPDF (fitz) is not installed. Cannot copy-paste between PDF files. Install it via 'pip install pymupdf'."}

        # 1. Extract text from source
        src_doc = fitz.open(source_file)
        src_target_pages = _parse_page_numbers(source_pages, len(src_doc))
        extracted_text = ""
        
        for page_num in src_target_pages:
            page = src_doc[page_num]
            text_instances = _fuzzy_search_pdf_page(page, source_search_text, flags=0)
            if text_instances:
                # To get the surrounding text block, we extract text from the page
                # and find the closest match to our search text.
                full_page_text = page.get_text("text")
                norm_page = _normalize_text_for_search(full_page_text)
                norm_search = _normalize_text_for_search(source_search_text)
                
                # If exact normalized match exists, extract a slightly larger context
                if norm_search in norm_page:
                    start_idx = norm_page.find(norm_search)
                    # Extract a bit before and after (e.g. 200 chars) to capture the full paragraph/block
                    start = max(0, start_idx - 100)
                    end = min(len(norm_page), start_idx + len(norm_search) + 500)
                    extracted_text = norm_page[start:end].strip()
                    break
        
        src_doc.close()
        
        if not extracted_text:
            return {"success": False, "error": f"Source text not found in {source_file}."}

        # 2. Insert into target
        tgt_doc = fitz.open(target_file)
        tgt_target_pages = _parse_page_numbers(target_pages, len(tgt_doc))
        insert_count = 0
        
        for page_num in tgt_target_pages:
            page = tgt_doc[page_num]
            anchor_instances = _fuzzy_search_pdf_page(page, target_anchor_text, flags=0)
            if anchor_instances:
                rect = anchor_instances[0]
                point = fitz.Point(rect.x1, rect.y0)
                # Insert with a line break for readability
                page.insert_text(point, "\n" + extracted_text + "\n", fontsize=11, fontname="helv")
                insert_count += 1
                break
                
        if insert_count == 0:
            tgt_doc.close()
            return {"success": False, "error": f"Anchor text not found in {target_file}."}
            
        tgt_doc.save(output_path)
        tgt_doc.close()
        
        return {
            "success": True,
            "output": f"Successfully copied text from {source_file} and pasted into {output_path}.",
            "file_name": output_path
        }

    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": f"Copy-paste failed: {str(e)}",
            "traceback": traceback.format_exc()
        }

def tool_delete_document_page(
    file_name: str,
    pages_to_delete: str
) -> Dict[str, Any]:
    """
    Deletes specific pages from a PDF or PPTX document.

    Args:
        file_name (str): The path to the document file (PDF or PPTX).
        pages_to_delete (str): Comma-separated list of pages to delete (1-indexed). E.g., "1, 3, 5-7".
    """
    err = _check_file_exists(file_name)
    if err:
        return err

    file_ext = Path(file_name).suffix.lower()
    output_path = _get_output_path(file_name, "_page_deleted")

    try:
        if file_ext == ".pdf":
            if fitz is None:
                return {"success": False, "error": "PyMuPDF (fitz) is not installed. Cannot delete PDF pages. Install it via 'pip install pymupdf'."}
            doc = fitz.open(file_name)
            total_pages = len(doc)
            pages_to_del = _parse_page_numbers(pages_to_delete, total_pages)
            
            if not pages_to_del:
                return {"success": False, "error": "No valid pages specified for deletion."}
                
            # Delete pages from the end to avoid index shifting
            for page_num in sorted(pages_to_del, reverse=True):
                doc.delete_page(page_num)
                
            doc.save(output_path)
            doc.close()
            
        elif file_ext == ".pptx":
            from pptx import Presentation
            import copy
            
            prs = Presentation(file_name)
            total_slides = len(prs.slides)
            pages_to_del = _parse_page_numbers(pages_to_delete, total_slides)
            
            if not pages_to_del:
                return {"success": False, "error": "No valid pages specified for deletion."}
                
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            for page_num in sorted(pages_to_del, reverse=True):
                xml_slides.remove(slides[page_num])
                
            prs.save(output_path)
        else:
            return {"success": False, "error": f"Unsupported file extension: {file_ext}"}

        return {
            "success": True,
            "output": f"Successfully deleted pages {pages_to_delete} from {file_name}. Saved to {output_path}",
            "file_name": output_path
        }

    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": f"Page deletion failed: {str(e)}",
            "traceback": traceback.format_exc()
        }

def tool_replace_text_globally(
    file_name: str,
    search_text: str,
    replacement_text: str,
    match_case: bool = False
) -> Dict[str, Any]:
    """
    Replaces all occurrences of a text string throughout the entire document.
    This is a brute-force replacement and does NOT use fuzzy matching. 
    Useful for global find-and-replace tasks (e.g., changing a name or terminology).

    Args:
        file_name (str): The path to the document file (PDF, DOCX, or PPTX).
        search_text (str): The exact text to find.
        replacement_text (str): The text to replace it with.
        match_case (bool, optional): Whether to match the exact case. Defaults to False.
    """
    err = _check_file_exists(file_name)
    if err:
        return err

    file_ext = Path(file_name).suffix.lower()
    output_path = _get_output_path(file_name, "_global_replace")
    total_replaced = 0

    try:
        if file_ext == ".pdf":
            if fitz is None:
                return {"success": False, "error": "PyMuPDF (fitz) is not installed. Cannot replace text in PDF files. Install it via 'pip install pymupdf'."}
            doc = fitz.open(file_name)

            flags = 0 if match_case else fitz.TEXT_DEHYPHENATE
            
            for page in doc:
                text_instances = page.search_for(search_text, flags=flags)
                if text_instances:
                    for inst in text_instances:
                        page.add_redact_annot(inst, text=replacement_text, fill=(1, 1, 1), fontsize=11)
                        total_replaced += 1
                    page.apply_redactions()
                    
            if total_replaced == 0:
                doc.close()
                return {"success": False, "error": "Text not found."}
                
            doc.save(output_path)
            doc.close()

        elif file_ext == ".docx":
            import docx
            doc = docx.Document(file_name)
            
            for paragraph in doc.paragraphs:
                if search_text in paragraph.text:
                    # To maintain formatting, we replace in runs if possible, 
                    # but for simplicity and global reach, we replace in the full text and rewrite the paragraph.
                    # This may lose complex inline formatting within that paragraph.
                    new_text = paragraph.text.replace(search_text, replacement_text) if match_case else paragraph.text.lower().replace(search_text.lower(), replacement_text)
                    if new_text != paragraph.text:
                        total_replaced += paragraph.text.count(search_text)
                        for run in paragraph.runs:
                            run.text = ""
                        paragraph.runs[0].text = new_text if paragraph.runs else paragraph.add_run(new_text).text
                        
            if total_replaced == 0:
                return {"success": False, "error": "Text not found."}
                
            doc.save(output_path)

        elif file_ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_name)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        full_text = "".join(run.text for run in paragraph.runs)
                        if search_text in full_text:
                            new_text = full_text.replace(search_text, replacement_text) if match_case else full_text.lower().replace(search_text.lower(), replacement_text)
                            if new_text != full_text:
                                total_replaced += full_text.count(search_text)
                                for run in paragraph.runs:
                                    run.text = ""
                                if paragraph.runs:
                                    paragraph.runs[0].text = new_text
                                    
            if total_replaced == 0:
                return {"success": False, "error": "Text not found."}
                
            prs.save(output_path)
        else:
            return {"success": False, "error": f"Unsupported file extension: {file_ext}"}

        return {
            "success": True,
            "output": f"Successfully replaced {total_replaced} instance(s) of the text. Saved to {output_path}",
            "file_name": output_path
        }

    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": f"Global replacement failed: {str(e)}",
            "traceback": traceback.format_exc()
        }