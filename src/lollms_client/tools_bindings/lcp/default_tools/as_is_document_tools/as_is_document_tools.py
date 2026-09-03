"""
as_is_document_tools.py
=======================
Tools for inspecting, reading, and selectively searching native document formats
(.pdf, .docx, .xlsx, .pptx, .txt, .csv, .json) in the workspace without prompt context bloat.
"""

import os
import re
import base64
from pathlib import Path
from typing import Dict, Any, List, Optional
from ascii_colors import ASCIIColors


TOOL_LIBRARY_NAME = "As-Is Document Tools"
TOOL_LIBRARY_DESC = "Inspect, search, and extract selective content from binary and structured workspace documents."
TOOL_LIBRARY_ICON = "📄"

_VISION_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp", ".gif"}


def init_tools_library() -> None:
    """Ensure optional parsing dependencies are available."""
    try:
        import pipmaster as pm
        pm.ensure_packages({
            "pypdf": ">=3.0.0",
            "python-docx": ">=0.8.11",
            "openpyxl": ">=3.0.0",
            "python-pptx": ">=0.6.21"
        })
    except Exception:
        pass


def _extract_pdf_images_for_vlm(pdf_path: Path, target_pages: List[int], max_images: int = 5) -> tuple[List[str], List[str]]:
    """
    Extracts embedded raster images from specified PDF pages as base64-encoded PNG strings.
    Uses PyMuPDF (fitz) if available; returns (images, media_types).
    """
    images_b64: List[str] = []
    media_types: List[str] = []

    try:
        import fitz
    except ImportError:
        return images_b64, media_types

    doc = None
    try:
        doc = fitz.open(str(pdf_path))
        for page_num in target_pages:
            if len(images_b64) >= max_images:
                break
            if page_num >= len(doc):
                continue
            page = doc[page_num]
            image_list = page.get_images(full=True)
            for img_info in image_list:
                if len(images_b64) >= max_images:
                    break
                xref = img_info[0]
                try:
                    base_img = doc.extract_image(xref)
                    img_bytes = base_img.get("image")
                    if not img_bytes:
                        continue
                    img_ext = base_img.get("ext", "png").lower()
                    if img_ext not in _VISION_IMAGE_EXTS:
                        img_ext = "png"
                    b64_str = base64.b64encode(img_bytes).decode("utf-8")
                    images_b64.append(b64_str)
                    media_types.append(f"image/{img_ext}")
                except Exception:
                    continue
    except Exception:
        pass
    finally:
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass

    return images_b64, media_types


def tool_inspect_document(file_name: str) -> Dict[str, Any]:
    """
    Inspects a document file to return structural metadata (format, page/sheet count, outline/headings).

    Args:
        file_name (str): Relative path to the document file in the workspace.
    """
    path = Path(file_name)
    if not path.exists():
        return {"success": False, "error": f"File '{file_name}' not found."}

    ext = path.suffix.lower()
    file_size = path.stat().st_size

    try:
        if ext == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                return {"success": False, "error": "pypdf is not installed. Please install it via 'pip install pypdf' to inspect PDF files."}
            try:
                reader = PdfReader(str(path))
                return {
                    "success": True,
                    "file_name": file_name,
                    "format": "PDF",
                    "size_bytes": file_size,
                    "pages_count": len(reader.pages),
                    "output": f"PDF '{file_name}': {len(reader.pages)} page(s), {file_size:,} bytes."
                }
            except Exception as pdf_err:
                return {"success": False, "error": f"Failed to read PDF '{file_name}': {pdf_err}. The file may be corrupted, encrypted, or require a password."}

        elif ext == ".docx":
            import docx
            doc = docx.Document(str(path))
            headings = [p.text.strip() for p in doc.paragraphs if p.style.name.startswith("Heading") and p.text.strip()]
            return {
                "success": True,
                "file_name": file_name,
                "format": "DOCX",
                "size_bytes": file_size,
                "paragraphs_count": len(doc.paragraphs),
                "headings": headings[:15],
                "output": f"DOCX '{file_name}': {len(doc.paragraphs)} paragraph(s), {len(headings)} heading(s)."
            }

        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True)
            sheets = wb.sheetnames
            wb.close()
            return {
                "success": True,
                "file_name": file_name,
                "format": "Excel",
                "size_bytes": file_size,
                "sheets": sheets,
                "output": f"Excel '{file_name}': {len(sheets)} sheet(s) ({', '.join(sheets)})."
            }

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            return {
                "success": True,
                "file_name": file_name,
                "format": "PPTX",
                "size_bytes": file_size,
                "slides_count": len(prs.slides),
                "output": f"PPTX '{file_name}': {len(prs.slides)} slide(s)."
            }

        else:
            text = path.read_text(encoding="utf-8", errors="ignore")
            lines = text.splitlines()
            return {
                "success": True,
                "file_name": file_name,
                "format": ext.upper().lstrip(".") or "Text",
                "size_bytes": file_size,
                "lines_count": len(lines),
                "output": f"Text file '{file_name}': {len(lines)} line(s), {file_size:,} bytes."
            }

    except Exception as e:
        import traceback as _tb
        return {"success": False, "error": f"Failed to inspect document '{file_name}': {e}", "traceback": _tb.format_exc()}


def tool_read_document_content(
    file_name: str,
    page_or_sheet: Optional[str] = None,
    max_chars: int = 8000
) -> Dict[str, Any]:
    """
    Extracts readable text content from a document (supports specific pages or sheets).

    Args:
        file_name (str): Relative path to the file in the workspace.
        page_or_sheet (str, optional): Specific page number (for PDF/PPTX) or sheet name (for Excel).
        max_chars (int, optional): Maximum characters to extract. Defaults to 8000.
    """
    if isinstance(max_chars, str):
        try:
            max_chars = int(max_chars)
        except ValueError:
            max_chars = 8000

    path = Path(file_name)
    if not path.exists():
        return {"success": False, "error": f"File '{file_name}' not found."}

    ext = path.suffix.lower()

    try:
        if ext == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                return {"success": False, "error": "pypdf is not installed. Please install it via 'pip install pypdf' to read PDF files."}
            try:
                reader = PdfReader(str(path))
                total_pages = len(reader.pages)
            except Exception as pdf_err:
                return {"success": False, "error": f"Failed to open PDF '{file_name}': {pdf_err}. The file may be corrupted, encrypted, or require a password."}

            target_pages = list(range(total_pages))
            if page_or_sheet:
                target_pages = []
                for part in str(page_or_sheet).split(','):
                    part = part.strip()
                    if '-' in part:
                        try:
                            start, end = map(int, part.split('-'))
                            target_pages.extend(range(max(1, start) - 1, min(total_pages, end)))
                        except ValueError:
                            pass
                    elif part.isdigit():
                        p = int(part) - 1
                        if 0 <= p < total_pages:
                            target_pages.append(p)

            if not target_pages:
                return {"success": False, "error": f"No valid pages found in range '{page_or_sheet}'. Document has {total_pages} pages."}

            pages_text = []
            current_chars = 0
            for page_num in target_pages:
                if current_chars >= max_chars:
                    break
                page = reader.pages[page_num]
                page_text = page.extract_text() or ""
                pages_text.append(f"--- Page {page_num + 1} ---\n{page_text}")
                current_chars += len(page_text)

            full_text = "\n\n".join(pages_text)
            if len(full_text) > max_chars:
                full_text = full_text[:max_chars] + f"\n\n... [truncated, {len(full_text) - max_chars} more chars]"

            images_b64: List[str] = []
            media_types: List[str] = []
            has_no_text = not full_text.strip()

            if has_no_text or "[truncated" in full_text:
                images_b64, media_types = _extract_pdf_images_for_vlm(path, target_pages, max_images=5)

            if has_no_text and not images_b64:
                return {"success": False, "error": f"Pages {page_or_sheet or '1-' + str(total_pages)} extracted but contained no readable text or extractable images. The PDF may require OCR."}
            if has_no_text and images_b64:
                return {
                    "success": True,
                    "total_pages": total_pages,
                    "pages_read": [p + 1 for p in target_pages],
                    "content": f"[Pages {page_or_sheet} contain no extractable text. {len(images_b64)} embedded figure(s)/image(s) were extracted for visual analysis.]",
                    "output": f"[Pages {page_or_sheet} contain no extractable text. {len(images_b64)} embedded figure(s)/image(s) were extracted for visual analysis.]",
                    "images": images_b64,
                    "image_media_types": media_types,
                    "prompt_injection": f"\n\n[SYSTEM: {len(images_b64)} image(s) from page(s) {page_or_sheet} have been extracted and attached as vision content for your analysis. Describe what you see in the figures, including any equations, charts, or diagrams.]"
                }

            result = {"success": True, "total_pages": total_pages, "pages_read": [p + 1 for p in target_pages], "content": full_text, "output": full_text}
            if images_b64:
                result["images"] = images_b64
                result["image_media_types"] = media_types
                result["prompt_injection"] = f"\n\n[SYSTEM: {len(images_b64)} embedded figure(s)/image(s) from page(s) {page_or_sheet} have been extracted and attached as vision content. Reference them in your analysis of the text.]"
            return result

        elif ext == ".docx":
            import docx
            doc = docx.Document(str(path))
            text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())[:max_chars]
            return {"success": True, "content": text, "output": text}

        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            sheet_name = page_or_sheet if page_or_sheet in wb.sheetnames else wb.sheetnames[0]
            ws = wb[sheet_name]

            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
                if sum(len(r) for r in rows) >= max_chars:
                    break

            wb.close()
            content = f"### Sheet: {sheet_name}\n\n" + "\n".join(rows)
            return {"success": True, "sheet": sheet_name, "content": content[:max_chars], "output": content[:max_chars]}

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                parts = [f"## Slide {i + 1}"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if p.text.strip():
                                parts.append(p.text.strip())
                slides_text.append("\n".join(parts))
                if sum(len(s) for s in slides_text) >= max_chars:
                    break

            content = "\n\n---\n\n".join(slides_text)[:max_chars]
            return {"success": True, "slides_count": len(prs.slides), "content": content, "output": content}

        else:
            text = path.read_text(encoding="utf-8", errors="ignore")[:max_chars]
            if not text.strip():
                return {"success": False, "error": f"File '{file_name}' exists but contains no readable text."}
            return {"success": True, "content": text, "output": text}

    except Exception as e:
        return {"success": False, "error": f"Failed to read content from '{file_name}': {e}"}


def tool_grep_document(
    file_name: str,
    pattern: str,
    max_matches: int = 20
) -> Dict[str, Any]:
    """
    Searches for a regular expression or keyword within a document without loading the whole file into context.

    Args:
        file_name (str): Relative path to the file in the workspace.
        pattern (str): Search keyword or regex pattern.
        max_matches (int, optional): Maximum matching excerpts to return. Defaults to 20.
    """
    path = Path(file_name)
    if not path.exists():
        return {"success": False, "error": f"File '{file_name}' not found."}

    regex = re.compile(pattern, re.IGNORECASE)
    matches = []

    try:
        doc_result = tool_read_document_content(file_name, max_chars=100000)
        if not doc_result.get("success"):
            return doc_result

        lines = doc_result.get("content", "").splitlines()
        for idx, line in enumerate(lines, 1):
            if regex.search(line):
                matches.append(f"Line {idx}: {line.strip()[:300]}")
                if len(matches) >= max_matches:
                    break

        output_str = f"Found {len(matches)} match(es) for '{pattern}' in '{file_name}':\n" + "\n".join(matches) if matches else f"No matches found for '{pattern}' in '{file_name}'."
        return {
            "success": True,
            "matches_count": len(matches),
            "matches": matches,
            "output": output_str
        }

    except Exception as e:
        return {"success": False, "error": f"Grep failed on '{file_name}': {e}"}