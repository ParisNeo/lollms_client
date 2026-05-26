# lollms_discussion/_mixin_export.py
from __future__ import annotations

import io
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, TYPE_CHECKING

from ascii_colors import ASCIIColors, trace_exception
from ._artefacts import ArtefactType

def _parse_html_slides(html_content: str) -> List[Dict[str, Any]]:
    import pipmaster as pm
    try:
        pm.ensure_packages("beautifulsoup4")
        from bs4 import BeautifulSoup
    except Exception as e:
        ASCIIColors.warning(f"Could not load beautifulsoup4: {e}. Fallback to regex-based slide parser.")
        slides = []
        sections = re.findall(r'<section\s+class="slide[^"]*"[^>]*>(.*?)</section>', html_content, re.DOTALL | re.IGNORECASE)
        for idx, sec in enumerate(sections):
            notes_match = re.search(r'data-notes="([^"]*)"', sec, re.IGNORECASE)
            notes = notes_match.group(1).strip() if notes_match else ""

            title_match = re.search(r'<h[12][^>]*>(.*?)</h[12]>', sec, re.DOTALL | re.IGNORECASE)
            title = re.sub(r'<[^>]+>', '', title_match.group(1)).strip() if title_match else ""

            subtitle_match = re.search(r'<h3[^>]*>(.*?)</h3>', sec, re.DOTALL | re.IGNORECASE)
            subtitle = re.sub(r'<[^>]+>', '', subtitle_match.group(1)).strip() if subtitle_match else ""

            bullets = []
            for li in re.findall(r'<li[^>]*>(.*?)</li>', sec, re.DOTALL | re.IGNORECASE):
                bullets.append(re.sub(r'<[^>]+>', '', li).strip())

            slides.append({
                "index": idx,
                "title": title,
                "subtitle": subtitle,
                "bullets": [b for b in bullets if b],
                "notes": notes
            })
        return slides

    soup = BeautifulSoup(html_content, "html.parser")
    slides = []

    sections = soup.find_all("section", class_="slide")
    if not sections:
        sections = soup.find_all("section")

    for idx, sec in enumerate(sections):
        notes = sec.get("data-notes", "").strip()

        title_el = sec.find(class_="slide-title") or sec.find(["h1", "h2"])
        title = title_el.get_text().strip() if title_el else ""

        subtitle_el = sec.find(class_="slide-subtitle") or sec.find("h3")
        subtitle = subtitle_el.get_text().strip() if subtitle_el else ""

        bullets = []
        bullet_list = sec.find(class_="slide-bullets") or sec.find(["ul", "ol"])
        if bullet_list:
            bullets = [li.get_text().strip() for li in bullet_list.find_all("li") if li.get_text().strip()]
        else:
            body_el = sec.find(class_="slide-body")
            if body_el:
                bullets = [p.get_text().strip() for p in body_el.find_all(["p", "div"]) if p.get_text().strip()]
            else:
                bullets = [p.get_text().strip() for p in sec.find_all("p") if p.get_text().strip()]

        if not title and (sec.find("h1") or sec.find("h2")):
            h = sec.find(["h1", "h2"])
            title = h.get_text().strip()

        slides.append({
            "index": idx,
            "title": title,
            "subtitle": subtitle,
            "bullets": bullets,
            "notes": notes
        })

    return slides

class ExportMixin:
    """Provides methods to export artifacts in various formats."""
    
    def export_artifact(
        self,
        title: str,
        export_format: str,
        version: Optional[int] = None
    ) -> Tuple[bytes, str]:
        """
        Exports an artifact of the given title in the requested format.
        
        Returns:
            Tuple[bytes, str]: (binary_content, mime_type)
        """
        art = self.artefacts.get(title, version)
        if not art:
            raise ValueError(f"Artifact '{title}' not found.")
            
        fmt = export_format.lower().strip()
        atype = art.get("type", "document")
        
        if atype == "data":
            # Data-based exports
            if fmt == "csv":
                return self._export_data_as_csv(art), "text/csv"
            elif fmt in ("excel", "xlsx"):
                return self._export_data_as_excel(art), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            else:
                raise ValueError(f"Data artifacts can only be exported as 'csv' or 'excel'.")
        else:
            # Text-based exports
            if fmt == "markdown":
                return art["content"].encode("utf-8"), "text/markdown"
            elif fmt == "html":
                if atype == "presentation":
                    return art["content"].encode("utf-8"), "text/html"
                return self._export_text_as_html(art), "text/html"
            elif fmt == "pdf":
                return self._export_text_as_pdf(art), "application/pdf"
            elif fmt == "docx":
                return self._export_text_as_docx(art), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            elif fmt == "pptx":
                return self._export_text_as_pptx(art), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            else:
                raise ValueError(f"Unsupported export format '{export_format}' for text artifacts.")

    def _export_text_as_html(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("markdown")
        import markdown
        
        content = art.get("content", "")
        body_html = markdown.markdown(content)
        
        # Add basic aesthetic styling to the HTML template
        html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{art['title']}</title>
<style>
    body {{
        font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
        line-height: 1.6;
        max-width: 800px;
        margin: 40px auto;
        padding: 0 20px;
        color: #333;
    }}
    h1, h2, h3 {{ color: #111; margin-top: 24px; }}
    pre {{ background: #f4f4f4; padding: 12px; border-radius: 6px; overflow-x: auto; }}
    code {{ font-family: monospace; background: #f4f4f4; padding: 2px 4px; border-radius: 3px; }}
    table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
    th, td {{ border: 1px solid #ddd; padding: 8px 12px; text-align: left; }}
    th {{ background: #f5f5f5; }}
</style>
</head>
<body>
{body_html}
</body>
</html>
"""
        return html.encode("utf-8")

    def _export_text_as_pdf(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("reportlab")
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors

        content = art.get("content", "")
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)

        styles = getSampleStyleSheet()
        story = []

        # Title
        title_style = ParagraphStyle(
            'TitleStyle',
            parent=styles['Heading1'],
            fontSize=24,
            leading=28,
            textColor=colors.HexColor('#1e293b'),
            spaceAfter=20
        )
        story.append(Paragraph(art["title"], title_style))

        atype = art.get("type", "document")
        if atype == "presentation" or "article class=\"presentation\"" in content:
            slides_data = _parse_html_slides(content)
            if slides_data:
                normal_style = styles['Normal']
                normal_style.fontSize = 12
                normal_style.leading = 16
                normal_style.spaceAfter = 10

                h1_style = ParagraphStyle('H1Style', parent=styles['Heading2'], fontSize=20, leading=24, spaceBefore=12, spaceAfter=8, textColor=colors.HexColor('#0f172a'))
                h2_style = ParagraphStyle('H2Style', parent=styles['Heading3'], fontSize=15, leading=18, spaceBefore=10, spaceAfter=6, textColor=colors.HexColor('#334155'))

                for s_idx, s_data in enumerate(slides_data):
                    if s_idx > 0:
                        story.append(PageBreak())
                    story.append(Paragraph(f"Slide {s_idx + 1}: {s_data['title']}", h1_style))
                    if s_data["subtitle"]:
                        story.append(Paragraph(s_data["subtitle"], h2_style))
                    story.append(Spacer(1, 10))
                    for bullet in s_data["bullets"]:
                        story.append(Paragraph(f"• {bullet}", normal_style))
                    if s_data["notes"]:
                        story.append(Spacer(1, 15))
                        notes_style = ParagraphStyle('NotesStyle', parent=normal_style, fontName='Helvetica-Oblique', fontSize=10, leading=13, textColor=colors.HexColor('#475569'))
                        story.append(Paragraph(f"<i>Speaker Notes:</i> {s_data['notes']}", notes_style))

                doc.build(story)
                return buffer.getvalue()

        # Parse simple markdown into Paragraphs
        normal_style = styles['Normal']
        normal_style.fontSize = 10
        normal_style.leading = 14
        normal_style.spaceAfter = 8

        h1_style = ParagraphStyle('H1Style', parent=styles['Heading2'], fontSize=16, leading=20, spaceBefore=12, spaceAfter=6, textColor=colors.HexColor('#0f172a'))
        h2_style = ParagraphStyle('H2Style', parent=styles['Heading3'], fontSize=13, leading=16, spaceBefore=10, spaceAfter=4, textColor=colors.HexColor('#334155'))

        code_style = ParagraphStyle('CodeStyle', parent=normal_style, fontName='Courier', fontSize=9, leading=11, spaceAfter=8, backColor=colors.HexColor('#f1f5f9'), borderPadding=6)

        lines = content.splitlines()
        in_code = False
        code_block = []

        for line in lines:
            line_stripped = line.strip()
            if line_stripped.startswith("```"):
                if in_code:
                    # End code block
                    code_text = "<br/>".join(code_block)
                    story.append(Paragraph(code_text, code_style))
                    code_block = []
                    in_code = False
                else:
                    in_code = True
                continue

            if in_code:
                # Escape HTML tags inside code
                escaped = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                code_block.append(escaped)
                continue

            if line_stripped.startswith("# "):
                story.append(Paragraph(line_stripped[2:], h1_style))
            elif line_stripped.startswith("## "):
                story.append(Paragraph(line_stripped[3:], h1_style))
            elif line_stripped.startswith("### "):
                story.append(Paragraph(line_stripped[4:], h2_style))
            elif line_stripped:
                # Basic inline bold/italic parser (simple)
                parsed_line = line_stripped
                parsed_line = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', parsed_line)
                parsed_line = re.sub(r'\*(.*?)\*', r'<i>\1</i>', parsed_line)
                story.append(Paragraph(parsed_line, normal_style))
            else:
                story.append(Spacer(1, 6))

        doc.build(story)
        return buffer.getvalue()

    def _export_text_as_docx(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("python-docx")
        import docx
        
        content = art.get("content", "")
        doc = docx.Document()
        doc.add_heading(art["title"], level=0)
        
        lines = content.splitlines()
        in_code = False
        code_p = None
        
        for line in lines:
            line_stripped = line.strip()
            if line_stripped.startswith("```"):
                if in_code:
                    in_code = False
                    code_p = None
                else:
                    in_code = True
                    code_p = doc.add_paragraph()
                    code_p.style = 'Normal'
                continue
                
            if in_code:
                if code_p:
                    run = code_p.add_run(line + "\n")
                    run.font.name = 'Courier New'
                continue
                
            if line_stripped.startswith("# "):
                doc.add_heading(line_stripped[2:], level=1)
            elif line_stripped.startswith("## "):
                doc.add_heading(line_stripped[3:], level=2)
            elif line_stripped.startswith("### "):
                doc.add_heading(line_stripped[4:], level=3)
            elif line_stripped:
                doc.add_paragraph(line_stripped)
                
        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    def _export_text_as_pptx(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("python-pptx")
        from pptx import Presentation
        from pptx.util import Inches, Pt

        content = art.get("content", "")
        prs = Presentation()

        atype = art.get("type", "document")
        if atype == "presentation" or "article class=\"presentation\"" in content:
            slides_data = _parse_html_slides(content)
            if slides_data:
                # First slide is title slide
                first = slides_data[0]
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = first["title"] or art["title"]
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = first["subtitle"] or "Exported from LoLLMS"
                if first["notes"] and slide.notes_slide:
                    slide.notes_slide.notes_text_frame.text = first["notes"]

                # Remaining slides
                slide_layout_content = prs.slide_layouts[1]
                for s_data in slides_data[1:]:
                    slide = prs.slides.add_slide(slide_layout_content)
                    slide.shapes.title.text = s_data["title"]
                    tf = slide.placeholders[1].text_frame
                    tf.text = s_data["subtitle"] if s_data["subtitle"] else ""
                    for bullet in s_data["bullets"]:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                    if s_data["notes"] and slide.notes_slide:
                        slide.notes_slide.notes_text_frame.text = s_data["notes"]

                buffer = io.BytesIO()
                prs.save(buffer)
                return buffer.getvalue()

        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = art["title"]
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Exported from LoLLMS Multimodal Viewer"

        # Split content into slides on h1/h2 headings
        lines = content.splitlines()
        slide_layout_content = prs.slide_layouts[1]
        current_slide = None
        tf = None

        for line in lines:
            line_stripped = line.strip()
            if line_stripped.startswith("# ") or line_stripped.startswith("## "):
                title_text = line_stripped.split(" ", 1)[1]
                current_slide = prs.slides.add_slide(slide_layout_content)
                current_slide.shapes.title.text = title_text
                tf = current_slide.placeholders[1].text_frame
            elif line_stripped:
                if tf:
                    p = tf.add_paragraph()
                    p.text = line_stripped
                    p.level = 0

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _export_data_as_csv(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages("pandas")
        import pandas as pd
        
        ext = art.get("file_ext", ".csv")
        title = art["title"]
        version = art.get("version", 1)
        workspace_dir = Path("./data_workspace")
        file_path = workspace_dir / f"{title}_v{version}{ext}"
        
        if not file_path.exists():
            raise FileNotFoundError(f"Original dataset file not found at {file_path}")
            
        try:
            if ext in (".db", ".sqlite", ".sqlite3"):
                import sqlite3
                conn = sqlite3.connect(str(file_path))
                cursor = conn.cursor()
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = [row[0] for row in cursor.fetchall()]
                if not tables:
                    raise ValueError("No tables found in SQLite DB.")
                df = pd.read_sql_query(f"SELECT * FROM {tables[0]};", conn)
                conn.close()
            elif ext in (".xlsx", ".xls"):
                df = pd.read_excel(str(file_path))
            else:
                sep = ";" if ext == ".csv" and ";" in file_path.read_text(encoding="utf-8", errors="ignore").splitlines()[0] else ","
                df = pd.read_csv(str(file_path), sep=sep)
                
            out = io.StringIO()
            df.to_csv(out, index=False)
            return out.getvalue().encode("utf-8")
        except Exception as e:
            trace_exception(e)
            raise RuntimeError(f"CSV export failed: {e}")

    def _export_data_as_excel(self, art: Dict[str, Any]) -> bytes:
        import pipmaster as pm
        pm.ensure_packages(["pandas", "openpyxl"])
        import pandas as pd
        
        ext = art.get("file_ext", ".csv")
        title = art["title"]
        version = art.get("version", 1)
        workspace_dir = Path("./data_workspace")
        file_path = workspace_dir / f"{title}_v{version}{ext}"
        
        if not file_path.exists():
            raise FileNotFoundError(f"Original dataset file not found at {file_path}")
            
        try:
            buffer = io.BytesIO()
            with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
                if ext in (".db", ".sqlite", ".sqlite3"):
                    import sqlite3
                    conn = sqlite3.connect(str(file_path))
                    cursor = conn.cursor()
                    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                    tables = [row[0] for row in cursor.fetchall()]
                    for table in tables:
                        df = pd.read_sql_query(f"SELECT * FROM {table};", conn)
                        df.to_excel(writer, sheet_name=table[:31], index=False)
                    conn.close()
                elif ext in (".xlsx", ".xls"):
                    xl = pd.ExcelFile(str(file_path))
                    for sheet in xl.sheet_names:
                        df = pd.read_excel(str(file_path), sheet_name=sheet)
                        df.to_excel(writer, sheet_name=sheet[:31], index=False)
                else:
                    sep = ";" if ext == ".csv" and ";" in file_path.read_text(encoding="utf-8", errors="ignore").splitlines()[0] else ","
                    df = pd.read_csv(str(file_path), sep=sep)
                    df.to_excel(writer, sheet_name="Data", index=False)
                    
            return buffer.getvalue()
        except Exception as e:
            trace_exception(e)
            raise RuntimeError(f"Excel export failed: {e}")
