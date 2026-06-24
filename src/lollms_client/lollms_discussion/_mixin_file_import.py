# lollms_discussion/_file_import.py
# FileImportMixin — import files of many formats into the artefact system.
#
# IMPORT MODES
# ------------
#   "text"        – Extract plain text only (all text-based formats + PDF).
#                   Result: one text/document artefact, no image artefacts.
#
#   "text_images" – Extract text AND embedded/rendered images.
#                   Images become a *separate*, deactivated image artefact
#                   (title "<filename>::images").  The text artefact embeds
#                   <artefact_image id="<title>::N" /> anchors so the LLM
#                   can correlate pixel data with prose.
#                   Supported: PDF (rendered pages), DOCX (embedded pics),
#                   PPTX (slide thumbnails + embedded pics).
#
#   "images_only" – Rasterise / collect images only, no text extraction.
#                   Each image becomes one page/frame of the image artefact.
#                   Useful for scanned PDFs, pure image files, and slides
#                   where text extraction would be noisy.
#
#   "ocr"         – Render each page to an image, pass it to the LLM vision
#                   API for OCR, stitch the per-page transcripts into a
#                   single text artefact. No image artefact is created.
#                   Requires lollmsClient.llm to support vision input.
#
# IMAGE ARTEFACT CONVENTION
# -------------------------
# Image artefacts created by this module are DEACTIVATED by default.
# They are referenced only through the <artefact_image> anchors in the
# companion text artefact.  The user (or the application UI) can activate
# them if they want the raw pixel data injected into the LLM context.
#
# SUPPORTED FORMATS
# -----------------
# Text-only         : .txt, .md, .csv, .json, .yaml, .yml, .xml, .html,
#                     .htm, .rst, .log, .py, .js, .ts, .java, .c, .cpp,
#                     .cs, .go, .rs, .rb, .php, .swift, .kt, .sh, .bat,
#                     .ps1, .sql, .r, .tex, and any unknown extension
#                     that decodes as UTF-8.
# Rich documents    : .pdf, .docx, .pptx, .xlsx, .odt (text + images)
# Pure images       : .png, .jpg, .jpeg, .bmp, .gif, .webp, .tiff, .tif
#
# DEPENDENCIES (all optional — auto-installed via pipmaster when needed)
# ----------------------------------------------------------------------
#   pipmaster         – Automatic package installer          [pip install pipmaster]
#   pymupdf   (fitz)  – PDF rendering & text extraction      [pip install pymupdf]
#   python-docx       – DOCX parsing                         [pip install python-docx]
#   python-pptx       – PPTX parsing                         [pip install python-pptx]
#   Pillow    (PIL)   – Image I/O & resizing                 [pip install pillow]
#   openpyxl          – XLSX text extraction                 [pip install openpyxl]
#   pdf2image         – PDF rendering fallback               [pip install pdf2image]
#   pypdf             – PDF text extraction fallback         [pip install pypdf]

from __future__ import annotations

import base64
import io
import json
import mimetypes
import os
import re
import uuid
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple, TYPE_CHECKING

from ascii_colors import ASCIIColors

from ._artefacts import ArtefactType, make_image_id
from ._data_files import _parse_data_file

if TYPE_CHECKING:
    from lollms_discussion import LollmsDiscussion

# ── pipmaster integration for optional dependencies ─────────────────────────
try:
    import pipmaster as pm
    _PM_AVAILABLE = True
except ImportError:
    _PM_AVAILABLE = False

def _ensure_installed(
    package_name: str,
    import_name: Optional[str] = None,
    package_version: Optional[str] = None,
) -> None:
    """
    Ensure an optional package is installed using pipmaster.
    Falls back to a plain import check if pipmaster is not available.
    """
    if _PM_AVAILABLE:
        if package_version:
            pm.ensure_packages({package_name: package_version})
        else:
            pm.ensure_packages(package_name)
    else:
        # Try a plain import; let the caller handle ImportError
        __import__(import_name or package_name)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

IMPORT_MODE_TEXT        = "text"
IMPORT_MODE_TEXT_IMAGES = "text_images"
IMPORT_MODE_TEXT_EMBEDDED_IMAGES = "text_embedded_images"
IMPORT_MODE_IMAGES_ONLY = "images_only"
IMPORT_MODE_OCR         = "ocr"
IMPORT_MODE_DATA        = "data"
IMPORT_MODE_DATA_BUNDLE = "data_bundle"
IMPORT_MODE_AUDIO_STT   = "audio_stt"

ALL_IMPORT_MODES = {
    IMPORT_MODE_TEXT,
    IMPORT_MODE_TEXT_IMAGES,
    IMPORT_MODE_TEXT_EMBEDDED_IMAGES,
    IMPORT_MODE_IMAGES_ONLY,
    IMPORT_MODE_OCR,
    IMPORT_MODE_DATA,
    IMPORT_MODE_DATA_BUNDLE,
    IMPORT_MODE_AUDIO_STT,
}

# Extensions treated as source code → ArtefactType.CODE
_CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".cc", ".cxx",
    ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".kts", ".sh", ".bash",
    ".zsh", ".fish", ".bat", ".ps1", ".sql", ".r", ".m", ".lua", ".ex", ".exs",
    ".erl", ".hs", ".clj", ".scala", ".dart", ".zig", ".nim", ".v",
}

# Extensions that are "plain text documents"
_TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv",
    ".json", ".yaml", ".yml", ".xml", ".html", ".htm", ".tex",
    ".ini", ".toml", ".cfg", ".conf", ".properties",
}

# Extensions for rich documents
_RICH_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".odt"}

# Pure image extensions
_IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp",
    ".tiff", ".tif", ".svg",
}

# Max image dimension for LLM context (pixels, longest side)
_MAX_LLM_IMAGE_DIM = 1024
# JPEG quality for rendered page thumbnails
_PAGE_JPEG_QUALITY  = 85
# Max pixels per PDF page when rendering (DPI ~ 150)
_PDF_RENDER_DPI     = 150


# ---------------------------------------------------------------------------
# Small utility helpers
# ---------------------------------------------------------------------------

def _b64_encode(data: bytes) -> str:
    return base64.b64encode(data).decode("utf-8")


def _resize_image_bytes(
    img_bytes: bytes,
    max_dim: int = _MAX_LLM_IMAGE_DIM,
    fmt: str = "JPEG",
    quality: int = _PDF_RENDER_DPI,
) -> Tuple[bytes, str]:
    """
    Resize image so its longest side ≤ max_dim.
    Returns (bytes, media_type).
    """
    _ensure_installed("Pillow", "PIL")
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(img_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        w, h = img.size
        if max(w, h) > max_dim:
            ratio = max_dim / max(w, h)
            img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format=fmt, quality=_PAGE_JPEG_QUALITY)
        return buf.getvalue(), f"image/{fmt.lower()}"
    except Exception as e:
        ASCIIColors.warning(f"[FileImport] Image resize failed: {e}")
        return img_bytes, "image/jpeg"


def _pil_image_to_b64(img: Any, fmt: str = "JPEG") -> Tuple[str, str]:
    """Convert a PIL image to (base64_str, media_type)."""
    _ensure_installed("Pillow", "PIL")
    buf = io.BytesIO()
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")
    img.save(buf, format=fmt, quality=_PAGE_JPEG_QUALITY)
    return _b64_encode(buf.getvalue()), f"image/{fmt.lower()}"


def _detect_artefact_type(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in (".csv", ".tsv", ".xlsx", ".xls", ".db", ".sqlite", ".sqlite3", ".sqlconn"):
        return ArtefactType.DATA
    if ext in _CODE_EXTENSIONS:
        return ArtefactType.CODE
    if ext in _IMAGE_EXTENSIONS:
        return ArtefactType.IMAGE
    return ArtefactType.DOCUMENT


def _parse_yaml_frontmatter(content: str) -> Tuple[Optional[Dict[str, Any]], str]:
    """
    Parses YAML frontmatter enclosed by triple-dashes (---) near the top of a file.
    Enforces strict checks to differentiate skill files from Hugging Face model/dataset README cards,
    while remaining robust to leading BOMs, whitespace, or injected page headers.
    """
    import re
    content_stripped = content.strip()

    # Search for frontmatter block anywhere within the first 1000 characters
    m = re.search(r'---(.*?)---', content_stripped[:1000], re.DOTALL)
    if not m:
        return None, content

    yaml_str = m.group(1).strip()
    start_idx = m.start()
    end_idx = m.end()
    body = content_stripped[:start_idx] + content_stripped[end_idx:]

    metadata = {}
    for line in yaml_str.splitlines():
        line = line.strip()
        if not line or ":" not in line:
            continue
        k, v = line.split(":", 1)
        k = k.strip().lower()
        v = v.strip().strip("'\"")
        metadata[k] = v

    # Enforce presence of essential skill schema keys (name and category)
    # Hugging Face cards never use category or author as top-level metadata keys.
    if "name" not in metadata or "category" not in metadata:
        return None, content

    return metadata, body.strip()


def _detect_language(path: Path) -> Optional[str]:
    ext = path.suffix.lower()
    _MAP = {
        ".py": "python", ".js": "javascript", ".ts": "typescript",
        ".jsx": "jsx", ".tsx": "tsx", ".java": "java",
        ".c": "c", ".cpp": "cpp", ".cc": "cpp", ".cxx": "cpp",
        ".cs": "csharp", ".go": "go", ".rs": "rust", ".rb": "ruby",
        ".php": "php", ".swift": "swift", ".kt": "kotlin",
        ".sh": "bash", ".bash": "bash", ".zsh": "zsh",
        ".bat": "batch", ".ps1": "powershell", ".sql": "sql",
        ".r": "r", ".lua": "lua", ".html": "html", ".htm": "html",
        ".xml": "xml", ".json": "json", ".yaml": "yaml", ".yml": "yaml",
        ".toml": "toml", ".md": "markdown", ".tex": "latex",
        ".csv": "csv",
    }
    return _MAP.get(ext)


# ---------------------------------------------------------------------------
# Per-format extraction functions
# ---------------------------------------------------------------------------

# ── Plain text / code ──────────────────────────────────────────────────────

def _extract_text_file(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_bytes().decode("utf-8", errors="replace")


# ── PDF ────────────────────────────────────────────────────────────────────

def _extract_pdf_text(path: Path) -> str:
    """Extract text from all PDF pages as Markdown, preserving tables."""
    _ensure_installed("pymupdf4llm")
    try:
        import pymupdf4llm

        # to_markdown returns one big string covering all pages
        # page_chunks=True returns a list of dicts, one per page
        chunks: list[dict] = pymupdf4llm.to_markdown(str(path), page_chunks=True)

        pages = []
        for i, chunk in enumerate(chunks):
            md = (chunk.get("text") or "").strip()
            pages.append(f"## Page {i + 1}\n\n{md}" if md else f"## Page {i + 1}\n\n[No text]")

        return "\n\n".join(pages)

    except ImportError:
        ASCIIColors.warning("[FileImport] pymupdf4llm not installed — falling back to pypdf")
        _ensure_installed("pypdf")
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages = []
            for i, page in enumerate(reader.pages):
                text = (page.extract_text() or "").strip()
                pages.append(
                    f"## Page {i + 1}\n\n{text}" if text else f"## Page {i + 1}\n\n[No text]"
                )
            return "\n\n".join(pages)
        except Exception as e2:
            raise RuntimeError(f"PDF text extraction failed: {e2}") from e2


def _dataframe_to_markdown(df: Any) -> str:
    """Convert a pandas DataFrame to a Markdown table."""
    try:
        import pandas as pd
        if not isinstance(df, pd.DataFrame):
            return ""
        # Replace newlines in cells with spaces to avoid breaking table format
        clean_df = df.copy()
        for col in clean_df.columns:
            clean_df[col] = clean_df[col].astype(str).str.replace("\n", " ").str.replace("\r", " ")
        return clean_df.to_markdown(index=False)
    except Exception:
        # Minimal fallback if pandas to_markdown fails
        lines = []
        # Header
        headers = [str(c) for c in df.columns]
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for _, row in df.iterrows():
            lines.append("| " + " | ".join(str(v) for v in row) + " |")
        return "\n".join(lines)


def _pdf_text_with_embedded_images(path: Path, art_title: str, progress_cb: Optional[Callable[[str], None]] = None) -> Tuple[str, List[Tuple[str, str]]]:
    """
    Extract PDF text as Markdown and extract raw embedded images.
    Returns (text_with_anchors, [(b64, media_type), ...]).
    """
    _ensure_installed("pymupdf", "fitz")
    try:
        import fitz
        doc    = fitz.open(str(path))
        pages_text: List[str] = []
        images: List[Tuple[str, str]] = []

        _has_md = hasattr(fitz.Page, "get_text") and "markdown" in getattr(fitz, "TEXT_FORMATS", set())

        img_count = 0
        for i, page in enumerate(doc):
            if progress_cb:
                try:
                    progress_cb(f"Extracting embedded images from page {i + 1}/{len(doc)}...")
                except Exception:
                    pass
            if _has_md:
                md = page.get_text("markdown").strip()
                text_blk = md if md else "[No selectable text on this page]"
            else:
                text_blk = (page.get_text() or "").strip() or "[No selectable text on this page]"

            # Find and extract raw embedded images on this page
            image_list = page.get_images(full=True)
            page_anchors = []
            for img_info in image_list:
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                img_ext = base_image["ext"]
                mt = f"image/{img_ext}"

                # Resize if extremely large to prevent context bloat
                resized, mt2 = _resize_image_bytes(img_bytes, fmt=img_ext.upper())
                images.append((_b64_encode(resized), mt2))

                anchor = f'<artefact_image id="{make_image_id(art_title, img_count)}" />'
                page_anchors.append(anchor)
                img_count += 1

            anchors_text = "\n\n".join(page_anchors) + "\n\n" if page_anchors else ""
            pages_text.append(f"## Page {i + 1}\n\n{anchors_text}{text_blk}")

        doc.close()
        return "\n\n".join(pages_text), images
    except Exception as e:
        ASCIIColors.warning(f"Failed to extract raw embedded images: {e}. Falling back to standard text extraction.")
        return _extract_pdf_text(path), []


def _extract_pdf_pages_as_images(
    path: Path,
    dpi: int = _PDF_RENDER_DPI,
    max_dim: int = _MAX_LLM_IMAGE_DIM,
    progress_cb: Optional[Callable[[str], None]] = None,
) -> List[Tuple[str, str]]:
    """
    Render each PDF page to a JPEG image.
    Returns list of (base64_str, media_type).
    """
    _ensure_installed("pymupdf4llm", "fitz")
    try:
        import fitz
        doc   = fitz.open(str(path))
        zoom  = dpi / 72.0
        mat   = fitz.Matrix(zoom, zoom)
        pages = []
        for i, page in enumerate(doc):
            if progress_cb:
                try:
                    progress_cb(f"Rendering page {i + 1}/{len(doc)}...")
                except Exception:
                    pass
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("jpeg")
            resized, mt = _resize_image_bytes(img_bytes, max_dim=max_dim, fmt="JPEG")
            pages.append((_b64_encode(resized), mt))
        doc.close()
        return pages
    except ImportError:
        _ensure_installed("pdf2image")
        try:
            from pdf2image import convert_from_path
            pil_pages = convert_from_path(str(path), dpi=dpi)
            pages = []
            for i, img in enumerate(pil_pages):
                if progress_cb:
                    try:
                        progress_cb(f"Rendering page {i + 1}/{len(pil_pages)}...")
                    except Exception:
                        pass
                b64, mt = _pil_image_to_b64(img)
                img_bytes = base64.b64decode(b64)
                resized, mt2 = _resize_image_bytes(img_bytes, max_dim=max_dim)
                pages.append((_b64_encode(resized), mt2))
            return pages
        except Exception as e2:
            raise RuntimeError(f"PDF page rendering failed: {e2}") from e2


def _pdf_text_with_image_anchors(path: Path, art_title: str, progress_cb: Optional[Callable[[str], None]] = None) -> Tuple[str, List[Tuple[str, str]]]:
    """
    Extract PDF text as Markdown and render pages as images simultaneously.
    Returns (text_with_anchors, [(b64, media_type), ...]).
    """
    _ensure_installed("pymupdf", "fitz")
    try:
        import fitz
        doc    = fitz.open(str(path))
        zoom   = _PDF_RENDER_DPI / 72.0
        mat    = fitz.Matrix(zoom, zoom)
        pages_text: List[str] = []
        images: List[Tuple[str, str]] = []

        _has_md = hasattr(fitz.Page, "get_text") and "markdown" in getattr(fitz, "TEXT_FORMATS", set())

        for i, page in enumerate(doc):
            if progress_cb:
                try:
                    progress_cb(f"Rendering page {i + 1}/{len(doc)}...")
                except Exception:
                    pass
            pix      = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("jpeg")
            resized, mt = _resize_image_bytes(img_bytes)
            images.append((_b64_encode(resized), mt))

            anchor   = f'<artefact_image id="{make_image_id(art_title, i)}" />'

            if _has_md:
                md = page.get_text("markdown").strip()
                text_blk = md if md else "[No selectable text on this page]"
            else:
                # Structured fallback with table detection
                text_blocks = []
                blocks = page.get_text("blocks")
                if blocks:
                    for b in blocks:
                        if len(b) >= 5:
                            txt = str(b[4]).strip()
                            if txt:
                                text_blocks.append(txt)

                tables_md = []
                try:
                    tabs = page.find_tables()
                    for tab in tabs.tables:
                        df = tab.to_pandas()
                        if df is not None and not df.empty:
                            tables_md.append(_dataframe_to_markdown(df))
                except Exception:
                    pass

                if tables_md:
                    combined = "\n\n".join(text_blocks) + "\n\n" + "\n\n".join(tables_md)
                else:
                    combined = "\n\n".join(text_blocks)

                text_blk = combined if combined else "[No selectable text on this page]"

            pages_text.append(f"## Page {i + 1}\n\n{anchor}\n\n{text_blk}")

        doc.close()
        return "\n\n".join(pages_text), images
    except ImportError:
        # Fallback: text only from pypdf, images from pdf2image
        _ensure_installed("pypdf")
        _ensure_installed("pdf2image")
        text   = _extract_pdf_text(path)
        images = _extract_pdf_pages_as_images(path)

        # Re-build text with anchors
        page_blocks = text.split("\n\n## Page ")
        rebuilt: List[str] = []
        for i, block in enumerate(page_blocks):
            anchor = f'<artefact_image id="{make_image_id(art_title, i)}" />'
            if i == 0:
                rebuilt.append(f"{block}\n\n{anchor}")
            else:
                rebuilt.append(f"## Page {block}\n\n{anchor}")
        return "\n\n".join(rebuilt), images


# ── DOCX ───────────────────────────────────────────────────────────────────

def _extract_docx_text(path: Path) -> str:
    _ensure_installed("python-docx", "docx")
    try:
        import docx as _docx
        doc   = _docx.Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                level = para.style.name.split()[-1]
                try:
                    hashes = "#" * int(level)
                except ValueError:
                    hashes = "##"
                parts.append(f"{hashes} {para.text}")
            else:
                parts.append(para.text)
        for tbl in doc.tables:
            rows = []
            for row in tbl.rows:
                rows.append(" | ".join(c.text.strip() for c in row.cells))
            parts.append("\n".join(rows))
        return "\n\n".join(p for p in parts if p.strip())
    except Exception as e:
        raise RuntimeError(f"DOCX text extraction failed: {e}") from e


def _extract_docx_embedded_images(path: Path) -> List[Tuple[str, str]]:
    """Extract embedded images from a DOCX file."""
    _ensure_installed("python-docx", "docx")
    try:
        import docx as _docx
        import zipfile
        images: List[Tuple[str, str]] = []
        with zipfile.ZipFile(str(path)) as zf:
            for name in zf.namelist():
                if name.startswith("word/media/") and not name.endswith("/"):
                    raw = zf.read(name)
                    ext = Path(name).suffix.lower()
                    mt  = mimetypes.types_map.get(ext, "image/jpeg")
                    resized, mt2 = _resize_image_bytes(raw, fmt="JPEG")
                    images.append((_b64_encode(resized), mt2))
        return images
    except Exception as e:
        ASCIIColors.warning(f"[FileImport] DOCX image extraction failed: {e}")
        return []


def _docx_text_with_image_anchors(path: Path, art_title: str) -> Tuple[str, List[Tuple[str, str]]]:
    _ensure_installed("python-docx", "docx")
    text   = _extract_docx_text(path)
    images = _extract_docx_embedded_images(path)
    if not images:
        return text, []

    # Append image anchors at the end of the document
    anchors = "\n\n".join(
        f'<artefact_image id="{make_image_id(art_title, i)}" />'
        for i in range(len(images))
    )
    return f"{text}\n\n---\n\n## Embedded Images\n\n{anchors}", images


# ── PPTX ───────────────────────────────────────────────────────────────────

def _extract_pptx_text(path: Path) -> str:
    _ensure_installed("python-pptx", "pptx")
    try:
        from pptx import Presentation
        prs   = Presentation(str(path))
        slides: List[str] = []
        for i, slide in enumerate(prs.slides):
            parts = [f"## Slide {i + 1}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            parts.append(t)
            slides.append("\n\n".join(parts))
        return "\n\n---\n\n".join(slides)
    except Exception as e:
        raise RuntimeError(f"PPTX text extraction failed: {e}") from e


def _extract_pptx_slide_images(
    path: Path,
    max_dim: int = _MAX_LLM_IMAGE_DIM,
) -> List[Tuple[str, str]]:
    """
    Render PPTX slides as images using python-pptx + Pillow.
    Falls back to extracting embedded image blobs.
    """
    _ensure_installed("python-pptx", "pptx")
    _ensure_installed("Pillow", "PIL")
    images: List[Tuple[str, str]] = []
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(str(path))

        for slide in prs.slides:
            # Try to collect embedded images from each slide
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_bytes = shape.image.blob
                        resized, mt = _resize_image_bytes(img_bytes, max_dim=max_dim)
                        images.append((_b64_encode(resized), mt))
                    except Exception:
                        pass
    except Exception as e:
        ASCIIColors.warning(f"[FileImport] PPTX image extraction: {e}")
    return images


def _pptx_text_with_image_anchors(path: Path, art_title: str) -> Tuple[str, List[Tuple[str, str]]]:
    _ensure_installed("python-pptx", "pptx")
    text   = _extract_pptx_text(path)
    images = _extract_pptx_slide_images(path)
    if not images:
        return text, []

    # Insert image anchors into each slide block
    slide_blocks = re.split(r'\n\n---\n\n', text)
    rebuilt: List[str] = []
    for i, block in enumerate(slide_blocks):
        if i < len(images):
            anchor = f'<artefact_image id="{make_image_id(art_title, i)}" />'
            rebuilt.append(f"{block}\n\n{anchor}")
        else:
            rebuilt.append(block)
    return "\n\n---\n\n".join(rebuilt), images


# ── XLSX ───────────────────────────────────────────────────────────────────

def _extract_xlsx_text(path: Path) -> str:
    _ensure_installed("openpyxl")
    try:
        import openpyxl
        wb    = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n\n".join(parts)
    except Exception as e:
        raise RuntimeError(f"XLSX text extraction failed: {e}") from e


# ── Pure images ────────────────────────────────────────────────────────────

import difflib

def _normalize_column_name(col_name: str) -> str:
    """
    Normalize a column name for schema comparison.
    - Convert to lowercase
    - Strip whitespace
    - Replace spaces and special chars with underscores
    - Collapse multiple underscores
    """
    import re
    normalized = str(col_name).lower().strip()
    normalized = re.sub(r'[\s\-\.\']+', '_', normalized)  # Replace spaces/dots/hyphens with _
    normalized = re.sub(r'_+', '_', normalized)  # Collapse multiple underscores
    return normalized

def _fuzzy_match_columns(source_cols: List[str], target_cols: List[str], threshold: float = 0.85) -> Dict[str, str]:
    """
    Match source columns to target columns using fuzzy string matching.
    Returns a mapping: {source_col: matched_target_col} for matches above threshold.
    Unmatched source columns are not included in the returned dict.
    """
    if not source_cols or not target_cols:
        return {}

    # Normalize target columns for comparison
    normalized_targets = [_normalize_column_name(c) for c in target_cols]

    mapping = {}
    for src_col in source_cols:
        src_norm = _normalize_column_name(src_col)

        # Find best match using difflib (fuzzy matching)
        matches = difflib.get_close_matches(
            src_norm, normalized_targets, n=1, cutoff=threshold
        )

        if matches:
            matched_target_raw = target_cols[normalized_targets.index(matches[0])]
            mapping[src_col] = matched_target_raw

    return mapping

def _load_image_file(path: Path) -> Tuple[str, str]:
    """Load an image file, resize it, return (base64, media_type)."""
    _ensure_installed("Pillow", "PIL")
    try:
        from PIL import Image
        img = Image.open(str(path))
        if img.mode not in ("RGB", "L", "RGBA"):
            img = img.convert("RGB")
        b64, mt = _pil_image_to_b64(img)
        raw = base64.b64decode(b64)
        resized, mt2 = _resize_image_bytes(raw, max_dim=_MAX_LLM_IMAGE_DIM)
        return _b64_encode(resized), mt2
    except Exception as e:
        raise RuntimeError(f"Image load failed: {e}") from e


# ── OCR via LLM vision ─────────────────────────────────────────────────────

def _ocr_images_with_llm(
    images_b64: List[str],
    lollms_client: Any,
    progress_callback: Optional[Callable[[str], None]] = None,
) -> str:
    """
    Pass each image to the LLM vision API and ask it to transcribe.
    Returns stitched text.
    """
    pages: List[str] = []
    for i, b64 in enumerate(images_b64):
        if progress_callback:
            progress_callback(f"OCR: transcribing page {i + 1}/{len(images_b64)}…")
        try:
            result = lollms_client.generate_text(
                prompt=(
                    "You are an OCR engine. Transcribe ALL text visible in this image "
                    "exactly as it appears, preserving paragraph breaks. "
                    "Output only the transcribed text — no commentary, no markdown "
                    "code fences, no preamble."
                ),
                images=[b64],
                n_predict=4096,
                temperature=0.0,
            )
            pages.append(f"## Page {i + 1}\n\n{result.strip()}")
        except Exception as e:
            ASCIIColors.warning(f"[FileImport] OCR page {i + 1} failed: {e}")
            pages.append(f"## Page {i + 1}\n\n[OCR failed: {e}]")
    return "\n\n".join(pages)


# ---------------------------------------------------------------------------
# Main dispatcher
# ---------------------------------------------------------------------------


class FileImportMixin:
    """
    Adds import_file() to LollmsDiscussion.

    Usage
    -----
    result = discussion.import_file(
        path        = "/path/to/document.pdf",
        mode        = "text_images",   # "text" | "text_images" | "images_only" | "ocr"
        title       = None,            # auto-derived from filename if None
        activate    = True,            # activate text artefact immediately
        progress_cb = my_callback,     # optional fn(str) for status updates
    )

    Returns
    -------
    {
        "text_artefact":   dict | None,   # created/updated text artefact
        "image_artefact":  dict | None,   # created image artefact (deactivated)
        "mode":            str,
        "page_count":      int,
        "image_count":     int,
        "warnings":        [str],
    }
    """

    def import_file(
        self,
        path:        str | Path,
        mode:        str  = IMPORT_MODE_TEXT_IMAGES,
        title:       Optional[str] = None,
        activate:    bool = True,
        progress_cb: Optional[Callable[[str], None]] = None,
        description: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Import a file into the artefact system.

        Parameters
        ----------
        path        : Path to the file on disk.
        mode        : One of "text", "text_images", "images_only", "ocr".
        title       : Artefact title (defaults to filename stem).
        activate    : Whether to activate the text artefact immediately.
        progress_cb : Optional callable receiving progress strings.

        Returns
        -------
        Dict with keys: text_artefact, image_artefact, mode,
                        page_count, image_count, warnings.
        """
        path     = Path(path)
        warnings: List[str] = []

        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        if mode not in ALL_IMPORT_MODES:
            raise ValueError(
                f"Unknown import mode '{mode}'. "
                f"Choose from: {sorted(ALL_IMPORT_MODES)}"
            )

        ext   = path.suffix.lower()
        title = title or path.stem

        def _progress(msg: str):
            ASCIIColors.info(f"[FileImport] {msg}")
            if progress_cb:
                try:
                    progress_cb(msg)
                except Exception:
                    pass

        _progress(f"Importing '{path.name}' in mode='{mode}'")

        text_artefact:  Optional[Dict] = None
        image_artefact: Optional[Dict] = None
        page_count  = 0
        image_count = 0

        # ── Determine file category ──────────────────────────────────────────
        is_plain_text = ext in _TEXT_EXTENSIONS or ext in _CODE_EXTENSIONS
        is_pdf        = ext == ".pdf"
        is_docx       = ext == ".docx"
        is_pptx       = ext == ".pptx"
        is_xlsx       = ext == ".xlsx"
        is_image_file = ext in _IMAGE_EXTENSIONS
        is_audio_file = ext in (".mp3", ".wav", ".ogg", ".flac", ".m4a", ".wma", ".aac")

        # ── data_bundle mode (Folder Ingestion) ──────────────────────────────
        # ── data_bundle mode (Folder Ingestion) ──────────────────────────────
        if mode == IMPORT_MODE_DATA_BUNDLE:
            if not path.is_dir():
                warnings.append(f"data_bundle mode requires a directory path, got: {path}")
                mode = IMPORT_MODE_TEXT
            else:
                _progress(f"Scanning folder for data files: {path}")
                supported_exts = {".csv", ".tsv", ".xlsx", ".xls"} # SQLite handled separately or skipped for now in fusion logic
                csv_xlsx_files = []
                
                # 1. SCAN PHASE: Find all relevant files and extract schemas
                for root, _, filenames in os.walk(path):
                    for fname in filenames:
                        f_path = Path(root) / fname
                        if f_path.suffix.lower() in supported_exts:
                            csv_xlsx_files.append(f_path)

                if not csv_xlsx_files:
                    warnings.append("No supported CSV/XLSX data files found in the folder.")
                    mode = IMPORT_MODE_TEXT
                else:
                    _progress(f"Found {len(csv_xlsx_files)} data files. Analyzing schemas...")
                    _ensure_installed("pandas")
                    _ensure_installed("openpyxl")
                    import pandas as pd
                    import sqlite3
                    from collections import defaultdict

                    workspace_dir = Path("./data_workspace")
                    try:
                        from lollms_client.app.server import APP_WORKSPACE_DIR as awd
                        if awd is not None:
                            workspace_dir = awd
                    except ImportError:
                        pass
                    workspace_dir.mkdir(exist_ok=True)

                    # Schema fingerprinting with normalized column names for better grouping
                    schema_groups = defaultdict(list)
                    schema_info_map = {}  # signature -> (normalized_columns, raw_columns_sample)

                    for f_path in csv_xlsx_files:
                        try:
                            ext = f_path.suffix.lower()
                            if ext in (".csv", ".tsv"):
                                sep = "\t" if ext == ".tsv" else ("," if ";" not in f_path.read_text(encoding="utf-8", errors="ignore").splitlines()[0] else ";")
                                df = pd.read_csv(str(f_path), nrows=1)  # Read header only
                            elif ext in (".xlsx", ".xls"):
                                xl = pd.ExcelFile(str(f_path))
                                if len(xl.sheet_names) == 1:
                                    df = pd.read_excel(str(f_path), sheet_name=xl.sheet_names[0], nrows=1)
                                else:
                                    df = pd.read_excel(str(f_path), sheet_name=xl.sheet_names[0], nrows=1)
                            else:
                                continue

                            # Create fingerprint using NORMALIZED column names for resilient grouping
                            # We use .dtypes.str[0] to get simple type kind ('i', 'f', 'O', etc.) which is more stable than full dtype names
                            normalized_signature = tuple(sorted(
                                (_normalize_column_name(str(c)), str(d.kind)) 
                                for c, d in zip(df.columns, df.dtypes)
                            ))

                            schema_groups[normalized_signature].append(f_path)
                            # Store info once per group (using first file as representative)
                            if normalized_signature not in schema_info_map:
                                schema_info_map[normalized_signature] = list(df.columns)  # Keep raw names for reference

                        except Exception as e:
                            warnings.append(f"Skipped {f_path.name} during schema analysis: {e}")

                    _progress(f"Found {len(schema_groups)} unique data structures. Fusing duplicates...")
                    
                    db_path = workspace_dir / f"{title}_consolidated.db"
                    conn = sqlite3.connect(str(db_path))
                    
                    tables_created = 0
                    total_rows_ingested = 0
                    table_details = []

                    for idx, (sig, file_list) in enumerate(schema_groups.items()):
                        try:
                            # Read ALL files in this group and concatenate them WITH source tracking
                            dfs_to_concat = []
                            sample_rows_info = []  # Store first few rows + filenames for LLM naming

                            # First pass: read all files to establish master column set with fuzzy matching
                            file_data_list = []
                            reference_columns = None

                            for f_path in file_list:
                                ext = f_path.suffix.lower()
                                if progress_cb: progress_cb(f"Reading {f_path.name} for schema group {idx}...")

                                if ext in (".csv", ".tsv"):
                                    sep = "\t" if ext == ".tsv" else ("," if ";" not in f_path.read_text(encoding="utf-8", errors="ignore").splitlines()[0] else ";")
                                    df = pd.read_csv(str(f_path), sep=sep)
                                elif ext in (".xlsx", ".xls"):
                                    xl = pd.ExcelFile(str(f_path))
                                    sheet_name = xl.sheet_names[0]
                                    df = pd.read_excel(str(f_path), sheet_name=sheet_name)

                                file_data_list.append((f_path, df))

                            # Establish reference columns from first file
                            if file_data_list:
                                reference_columns = list(file_data_list[0][1].columns)

                                # Second pass: align all files to reference schema using fuzzy matching
                                for f_path, df in file_data_list:
                                    if progress_cb: progress_cb(f"Aligning {f_path.name} columns...")

                                    # Fuzzy match current file's columns to reference columns
                                    column_mapping = _fuzzy_match_columns(df.columns.tolist(), reference_columns)

                                    # Create a new dataframe with aligned columns
                                    aligned_df = df.copy()

                                    # Rename matched columns to reference names
                                    for src_col, target_col in column_mapping.items():
                                        if src_col != target_col:  # Only rename if different
                                            aligned_df.rename(columns={src_col: target_col}, inplace=True)

                                    # Add missing columns from reference that don't exist in this file
                                    existing_cols = set(aligned_df.columns)
                                    for ref_col in reference_columns:
                                        if ref_col not in existing_cols:
                                            aligned_df[ref_col] = None  # Fill missing with NaN

                                    # Remove extra columns not in reference (they shouldn't be there if grouped correctly, but safety check)
                                    aligned_df = aligned_df[[c for c in reference_columns if c in aligned_df.columns]]

                                    # Add source_file column AFTER alignment
                                    aligned_df['source_file'] = f_path.name
                                    dfs_to_concat.append(aligned_df)

                                    # Collect sample data for LLM naming (first 3 rows as representative samples)
                                    if len(sample_rows_info) < 5:
                                        sample_rows_info.append({
                                            "filename": f_path.name,
                                            "sample": aligned_df.head(3).to_dict('records')
                                        })

                            # Fuse into one massive dataframe per schema group
                            if dfs_to_concat:
                                fused_df = pd.concat(dfs_to_concat, ignore_index=True)

                                # Clean column names to ensure SQL compatibility (replace spaces/special chars with underscore)
                                fused_df.columns = [str(c).replace(" ", "_").replace("-", "_").replace(".", "_") for c in fused_df.columns]

                                # Generate meaningful table name using LLM
                                _ensure_installed("pipmaster")
                                import pipmaster as pm

                                # Prepare context for LLM naming
                                columns_info = "\n".join([f"  • {col}" for col in fused_df.columns])
                                sample_preview = f"\nSample data from {len(sample_rows_info)} source file(s):\n"
                                for sample_info in sample_rows_info[:3]:
                                    sample_preview += f"\n--- From: {sample_info['filename']} ---\n"
                                    for row_idx, row_data in enumerate(sample_info['sample'][:2]):
                                        preview = ", ".join([f"{k}: {v}" for k, v in list(row_data.items())[:5]])
                                        sample_preview += f"  Row {row_idx+1}: {preview}...\n"

                                naming_prompt = (
                                    "You are a data architect. Based on the following database schema and sample data, "
                                    "propose a meaningful table name that describes what this data represents.\n\n"
                                    "COLUMNS:\n" + columns_info + "\n\n" + sample_preview + "\n\n"
                                    "Return ONLY the table name (no explanations, no quotes). "
                                    "Use snake_case format. Keep it under 50 characters."
                                )

                                try:
                                    # Use LLM to generate table name if available
                                    llm_available = False
                                    llm_table_name = None

                                    if hasattr(self, 'lollmsClient') and self.lollmsClient:
                                        try:
                                            response = self.lollmsClient.generate_text(
                                                prompt=naming_prompt,
                                                n_predict=50,
                                                temperature=0.3
                                            )
                                            llm_table_name = response.strip().lower()[:50]
                                            if llm_table_name and not any(c in llm_table_name for c in '!"#$%&()*+,:;<=>?@[]^`{|}'):
                                                llm_available = True
                                        except Exception:
                                            pass

                                    # Use LLM name or fallback to generated name
                                    if llm_available and llm_table_name:
                                        table_name = f"data_{llm_table_name}"
                                    else:
                                        # Fallback: use first 3 words from first filename + columns signature
                                        base_name = '_'.join(file_list[0].stem.split('_')[:2])
                                        key_cols = fused_df.columns[:3] if len(fused_df.columns) >= 3 else fused_df.columns
                                        table_name = f"data_{base_name}_{'_'.join(key_cols)}"

                                    # Ensure valid SQL identifier
                                    table_name = re.sub(r'[^a-zA-Z0-9_]', '_', table_name.lower())
                                    if len(table_name) > 60: 
                                        table_name = table_name[:58] + "_"

                                except Exception as naming_err:
                                    # Fallback on error
                                    base_name = '_'.join(file_list[0].stem.split('_')[:2])
                                    key_cols = fused_df.columns[:3] if len(fused_df.columns) >= 3 else fused_df.columns
                                    table_name = f"data_{base_name}_{'_'.join(key_cols)}"

                                # Save to DB - sanitize table name for SQLite compatibility
                                safe_table_name = re.sub(r'[^a-zA-Z0-9_]', '_', table_name)
                                if len(safe_table_name) > 63: 
                                    safe_table_name = safe_table_name[:61] + "_"

                                # Save with sanitized name but keep original in metadata for display
                                fused_df.to_sql(safe_table_name, conn, if_exists='replace', index=False)

                                # Commit immediately to ensure data is persisted
                                conn.commit()

                                tables_created += 1
                                total_rows_ingested += len(fused_df)
                                table_details.append({
                                    'safe_name': safe_table_name,
                                    'display_name': table_name,
                                    'files_merged': len(file_list),
                                    'rows': len(fused_df)
                                })

                        except Exception as e:
                            warnings.append(f"Failed to fuse group {idx}: {e}")

                    conn.close()

                    # Build schema summary AFTER all tables are created with actual counts
                    schema_summary = [f"# Data Bundle: {title}\n"]
                    schema_summary.append(f"Format: Consolidated SQLite Database (Schema Fused)\n")
                    schema_summary.append(f"Source Folder: {path}\n")
                    schema_summary.append(f"Total Files Scanned: {len(csv_xlsx_files)}\n")
                    schema_summary.append(f"Unique Structures Found: {len(schema_groups)}\n\n")

                    if table_details:
                        for td in table_details:
                            schema_summary.append(
                                f"  • Table '{td['safe_name']}' (display: {td['display_name']}) "
                                f"({td['files_merged']} files merged | {td['rows']:,} rows)"
                            )

                    schema_summary.append(f"\nTotal Tables Created: {tables_created}")
                    schema_summary.append(f"Total Rows Ingested: {total_rows_ingested:,}")

                    # Append description if provided
                    if description and description.strip():
                        schema_summary.append("\n---\n")
                        schema_summary.append("## Description\n")
                        schema_summary.append(description.strip())

                    text = "\n".join(schema_summary)

                    # Create the artifact
                    # Store description in metadata for future reference
                    extra_meta = {}
                    if description and description.strip():
                        extra_meta["description"] = description.strip()

                    art = self.artefacts.add(
                        title=title,
                        artefact_type=ArtefactType.DATA,
                        content=text,
                        active=activate,
                        file_ext=".db",
                        version=1,
                        read_only=True,
                        **extra_meta
                    )
                    _progress(f"Data bundle consolidation complete. Reduced {len(csv_xlsx_files)} files to {tables_created} tables.")
                    return {
                        "text_artefact": art,
                        "image_artefact": None,
                        "mode": mode,
                        "page_count": 0,
                        "image_count": 0,
                        "warnings": warnings,
                    }

            # ── data mode (Single File) ──────────────────────────────────────────
            if mode == IMPORT_MODE_DATA:
                _progress("Analyzing structured data file...")

                # 🛡️ TITLE NORMALIZATION PROTOCOL
                # Ensure 'title' does not already contain the extension to prevent "file.csv.csv" duplicates.
                # We strip the extension if it matches the file's actual extension.
                clean_title = title
                if clean_title.lower().endswith(ext) and len(clean_title) > len(ext):
                    clean_title = clean_title[:-len(ext)]
                    _progress(f"Normalized title '{title}' -> '{clean_title}' to prevent extension duplication.")

                # 1. Parse and extract the Markdown schema description
                schema_text, images_data = _parse_data_file(path, clean_title, version=1, progress_cb=progress_cb)

                # 2. Extract the actual raw content (verbatim rows) for CSV/TSV
                if ext in (".csv", ".tsv"):
                    raw_content = _extract_text_file(path)
                else:
                    raw_content = schema_text  # Fallback for binary databases/spreadsheets

                # Append description if provided (before the schema content)
                full_description = schema_text
                if description and description.strip():
                    full_description = f"## Description\n{description.strip()}\n\n{schema_text}"

                atype = ArtefactType.DATA

                # Use clean_title for lookup and registration
                existing = self.artefacts.get(clean_title)
                if existing is None:
                    art = self.artefacts.add(
                        title=clean_title,
                        artefact_type=atype,
                        content=raw_content,
                        active=activate,
                        file_ext=ext,
                        version=1,
                        read_only=True, # Default newly imported data files to read-only
                        description=full_description
                    )
                else:
                    art = self.artefacts.update(
                        title=clean_title,
                        new_content=raw_content,
                        new_type=atype,
                        active=activate,
                        file_ext=ext,
                        version=1,
                        read_only=True, # Default newly imported data files to read-only
                        description=full_description
                    )
                _progress("Data analysis complete.")

                # 🛑 CRITICAL EXIT: Prevent fall-through to generic text processing
                return {
                    "text_artefact": art,
                    "image_artefact": None,
                    "mode": mode,
                    "page_count": 0,
                    "image_count": 0,
                    "warnings": warnings,
                }

        # ── images_only: pure image files ────────────────────────────────────
        if mode == IMPORT_MODE_IMAGES_ONLY:
            images_data: List[Tuple[str, str]] = []

            if is_image_file:
                _progress("Loading image file…")
                b64, mt = _load_image_file(path)
                images_data = [(b64, mt)]
                page_count  = 1

            elif is_pdf:
                _progress("Rendering PDF pages…")
                images_data = _extract_pdf_pages_as_images(path, progress_cb=progress_cb)
                page_count  = len(images_data)

            elif is_pptx:
                _progress("Extracting PPTX slide images…")
                images_data = _extract_pptx_slide_images(path)
                page_count  = len(images_data)

            elif is_docx:
                _progress("Extracting DOCX embedded images…")
                images_data = _extract_docx_embedded_images(path)
                page_count  = len(images_data)

            else:
                warnings.append(f"images_only mode: no image extraction for '{ext}' files.")

            if images_data:
                image_count = len(images_data)
                imgs_b64  = [x[0] for x in images_data]
                imgs_mt   = [x[1] for x in images_data]
                image_artefact = self.artefacts.add(
                    title             = f"{title}::images",
                    artefact_type     = ArtefactType.IMAGE,
                    content           = f"Images extracted from '{path.name}' ({image_count} page(s)).",
                    images            = imgs_b64,
                    image_media_types = imgs_mt,
                    active            = activate,   # Respect active parameter!
                )
                _progress(f"Image artefact created: '{title}::images' ({image_count} image(s), active={activate})")

                # Generate anchors text
                anchors_content = f"# {title} (Images Only)\n\n" + "\n\n".join(
                    f"## Page {i + 1}\n\n<artefact_image id=\"{make_image_id(f'{title}::images', i)}\" />"
                    for i in range(image_count)
                )
                text_artefact = self._import_save_text(title, anchors_content, path, activate)
                _progress(f"Companion text/anchors artefact created: '{title}' (active={activate})")

            return {
                "text_artefact":  text_artefact,
                "image_artefact": image_artefact,
                "mode":           mode,
                "page_count":     page_count,
                "image_count":    image_count,
                "warnings":       warnings,
            }

        # ── ocr mode ────────────────────────────────────────────────────────
        if mode == IMPORT_MODE_OCR:
            lc = getattr(self, "lollmsClient", None)
            if lc is None:
                raise RuntimeError("OCR mode requires lollmsClient to be set on the discussion.")

            images_data = []

            if is_image_file:
                _progress("Loading image for OCR…")
                b64, mt = _load_image_file(path)
                images_data = [(b64, mt)]

            elif is_pdf:
                _progress("Rendering PDF pages for OCR…")
                images_data = _extract_pdf_pages_as_images(path)

            elif is_pptx:
                _progress("Extracting PPTX images for OCR…")
                images_data = _extract_pptx_slide_images(path)
                if not images_data:
                    warnings.append("No slide images found in PPTX; falling back to text extraction.")
                    text = _extract_pptx_text(path)
                    text_artefact = self._import_save_text(title, text, path, activate)
                    return {
                        "text_artefact": text_artefact, "image_artefact": None,
                        "mode": mode, "page_count": 0, "image_count": 0, "warnings": warnings,
                    }

            elif is_docx:
                _progress("Extracting DOCX images for OCR…")
                images_data = _extract_docx_embedded_images(path)
                if not images_data:
                    warnings.append("No embedded images in DOCX; falling back to text extraction.")
                    text = _extract_docx_text(path)
                    text_artefact = self._import_save_text(title, text, path, activate)
                    return {
                        "text_artefact": text_artefact, "image_artefact": None,
                        "mode": mode, "page_count": 0, "image_count": 0, "warnings": warnings,
                    }

            else:
                warnings.append(f"OCR mode not applicable for '{ext}'; using text extraction.")
                text = _extract_text_file(path)
                text_artefact = self._import_save_text(title, text, path, activate)
                return {
                    "text_artefact": text_artefact, "image_artefact": None,
                    "mode": mode, "page_count": 0, "image_count": 0, "warnings": warnings,
                }

            page_count  = len(images_data)
            images_b64  = [x[0] for x in images_data]
            ocr_text    = _ocr_images_with_llm(images_b64, lc, _progress)
            _progress(f"OCR complete — {page_count} page(s) transcribed")

            text_artefact = self._import_save_text(title, ocr_text, path, activate)

            return {
                "text_artefact":  text_artefact,
                "image_artefact": None,
                "mode":           mode,
                "page_count":     page_count,
                "image_count":    0,
                "warnings":       warnings,
            }

        # ── text and text_images modes ───────────────────────────────────────
        include_images = (mode in (IMPORT_MODE_TEXT_IMAGES, IMPORT_MODE_TEXT_EMBEDDED_IMAGES))
        extract_embedded = (mode == IMPORT_MODE_TEXT_EMBEDDED_IMAGES)
        text:   str                    = ""
        images_data: List[Tuple[str, str]] = []

        if is_plain_text:
            _progress("Reading text file…")
            text = _extract_text_file(path)
            # No image extraction for plain text

        elif is_audio_file or mode == IMPORT_MODE_AUDIO_STT:
            _progress("Transcribing audio file via STT...")
            lc = getattr(self, "lollmsClient", None)
            if lc is None or not getattr(lc, "stt", None):
                raise RuntimeError("Audio file import requires a Speech-to-Text (STT) model to be configured and active in your settings.")
            try:
                raw_transcript = lc.transcribe_audio(path)
                text = f"# Audio Transcript: {title}\n*File source: {path.name}*\n\n{raw_transcript}"
                _progress("Audio transcription complete!")
            except Exception as e:
                raise RuntimeError(f"Audio transcription failed: {e}")

        elif is_pdf:
            if extract_embedded:
                _progress("Extracting PDF text + raw embedded images…")
                text, images_data = _pdf_text_with_embedded_images(path, title, progress_cb=progress_cb)
            elif include_images:
                _progress("Extracting PDF text + rendering pages…")
                text, images_data = _pdf_text_with_image_anchors(path, title, progress_cb=progress_cb)
            else:
                _progress("Extracting PDF text…")
                text = _extract_pdf_text(path)
            page_count = text.count("\n## Page ")

        elif is_docx:
            if include_images:
                _progress("Extracting DOCX text + embedded images…")
                text, images_data = _docx_text_with_image_anchors(path, title)
            else:
                _progress("Extracting DOCX text…")
                text = _extract_docx_text(path)

        elif is_pptx:
            if include_images:
                _progress("Extracting PPTX text + slide images…")
                text, images_data = _pptx_text_with_image_anchors(path, title)
            else:
                _progress("Extracting PPTX text…")
                text = _extract_pptx_text(path)

        elif is_xlsx:
            _progress("Extracting XLSX content…")
            text = _extract_xlsx_text(path)
            if include_images:
                warnings.append("XLSX files have no embedded images to extract.")

        elif is_image_file:
            # Single image: make it a text artefact with one anchor + image
            if include_images:
                _progress("Loading image…")
                b64, mt = _load_image_file(path)
                images_data = [(b64, mt)]
                anchor = f'<artefact_image id="{make_image_id(title, 0)}" />'
                text   = f"# {title}\n\n{anchor}\n"
            else:
                _progress("Image file in text mode — storing metadata only.")
                text = f"# {title}\n\n[Image file: {path.name}]"

        else:
            # Unknown extension — attempt UTF-8 text read
            try:
                _progress(f"Unknown extension '{ext}' — attempting text read…")
                text = _extract_text_file(path)
                warnings.append(f"Extension '{ext}' not natively supported; treated as plain text.")
            except Exception as e:
                raise RuntimeError(
                    f"Cannot import '{path.name}': unsupported format and binary read failed. {e}"
                )

        # ── Save text artefact ───────────────────────────────────────────────
        text_artefact = self._import_save_text(title, text, path, activate)

        # ── Save image artefact (deactivated) ────────────────────────────────
        if include_images and images_data:
            image_count = len(images_data)
            imgs_b64    = [x[0] for x in images_data]
            imgs_mt     = [x[1] for x in images_data]
            image_artefact = self.artefacts.add(
                title             = f"{title}::images",
                artefact_type     = ArtefactType.IMAGE,
                content           = (
                    f"Images extracted from '{path.name}' ({image_count} image(s)).\n"
                    f"Referenced from artefact '{title}' via <artefact_image> anchors."
                ),
                images            = imgs_b64,
                image_media_types = imgs_mt,
                active            = False,   # deactivated — injected via anchors only
            )
            _progress(
                f"Image artefact '{title}::images' created "
                f"({image_count} image(s), deactivated)"
            )

        _progress("Import complete.")
        return {
            "text_artefact":  text_artefact,
            "image_artefact": image_artefact,
            "mode":           mode,
            "page_count":     page_count or (len(images_data) if images_data else 0),
            "image_count":    image_count,
            "warnings":       warnings,
        }

    # ── Internal helpers ─────────────────────────────────────────────────────

    def _import_save_text(
        self,
        title:    str,
        text:     str,
        path:     Path,
        activate: bool,
    ) -> Dict:
        """Create or update the text artefact for the imported content."""
        # Parse YAML frontmatter dynamically to detect skills
        metadata, body = _parse_yaml_frontmatter(text)

        if metadata:
            title = metadata.get("name", title)
            atype = ArtefactType.SKILL
            text = body
            extra_data = {
                "description": metadata.get("description", ""),
                "category": metadata.get("category", "lollms_client/general"),
                "author": metadata.get("author", "Unknown"),
                "skill_version": metadata.get("version", "1.0.0"),
                "created_at": metadata.get("created", datetime.utcnow().isoformat())
            }
        else:
            atype = _detect_artefact_type(path)
            extra_data = {}
            if atype == ArtefactType.DATA:
                extra_data["file_ext"] = path.suffix.lower()

        language = _detect_language(path)

        existing = self.artefacts.get(title)
        if existing is None:
            art = self.artefacts.add(
                title         = title,
                artefact_type = atype,
                content       = text,
                language      = language,
                active        = activate,
                **extra_data
            )
            ASCIIColors.success(f"[FileImport] Created artefact '{title}' of type '{atype}'")
        else:
            art = self.artefacts.update(
                title       = title,
                new_content = text,
                new_type    = atype,
                language    = language,
                active      = activate,
                **extra_data
            )
            ASCIIColors.success(
                f"[FileImport] Updated artefact '{title}' → v{art.get('version', '?')} of type '{atype}'"
            )
        return art
